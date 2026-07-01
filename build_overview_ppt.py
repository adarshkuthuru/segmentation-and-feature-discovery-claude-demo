"""
Build the 2-slide executive overview, styled to match the Zenon 2026 template
(templates/Zenon_2026_Template.pptx): 10x5.625in canvas, navy/gold palette,
white cards with a gold left accent bar, Calibri throughout.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

TEMPLATE = "templates/Zenon_2026_Template.pptx"
OUT = "outputs/Segment_Discovery_Overview_2026-07-01.pptx"
FONT = "Calibri"

# ── Zenon brand palette (sampled directly from the template) ────────────────
NAVY  = RGBColor(0x1D, 0x2D, 0x44)   # headings, body text, chip fills
GOLD  = RGBColor(0xC8, 0xA1, 0x5A)   # accent bars, numbered chips, dividers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x71, 0x80, 0x96)   # footer / caption text only
BG    = RGBColor(0xF5, 0xF5, 0xF5)   # page background (matches template slide 2)


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


prs = Presentation(TEMPLATE)
while len(prs.slides._sldIdLst) > 0:
    delete_slide(prs, 0)
LAYOUT = prs.slide_layouts[0]  # DEFAULT


# ── helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=10, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = FONT
    return tf

def bullets(slide, lines, l, t, w, h, size=7.5, color=NAVY, leading=1.0):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = f"•  {line}"
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tf

def page_header(slide, title, subtitle):
    txt(slide, title, 0.35, 0.08, 9.3, 0.32, size=16, bold=True, color=NAVY)
    box(slide, 0.35, 0.42, 1.0, 0.035, GOLD)
    txt(slide, subtitle, 0.35, 0.48, 9.3, 0.28, size=9, color=GRAY)

def section_label(slide, text, l, t, w):
    txt(slide, text, l, t, w, 0.22, size=9.5, bold=True, color=NAVY)

def footer(slide, page_num):
    txt(slide, "© 2026 Zenon, LLC. All rights reserved.",
        0.35, 5.25, 4.0, 0.25, size=8, color=GRAY)
    txt(slide, str(page_num), 9.3, 5.25, 0.4, 0.25, size=8, color=GRAY,
        align=PP_ALIGN.RIGHT)

def gold_bar_card(slide, l, t, w, h):
    box(slide, l, t, w, h, WHITE)
    box(slide, l, t, 0.05, h, GOLD)

def callout(slide, text, l=0.35, t=4.78, w=9.3, h=0.38, size=9.5):
    gold_bar_card(slide, l, t, w, h)
    txt(slide, text, l + 0.22, t, w - 0.4, h, size=size, bold=True,
        color=NAVY, align=PP_ALIGN.LEFT)
    # vertically center the single line
    tb = slide.shapes[-1]
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    tb.top = Inches(t + (h - Pt(size).inches * 1.3) / 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — SOLUTION ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(LAYOUT)
bg(s1)
page_header(s1,
    "Solution Architecture: What's Generic vs What's Yours",
    "Generic components work on any dataset — example-specific facts live in the config file or the prompt")
footer(s1, 1)

CONTENT_TOP = 0.80
CONTENT_BOT = 4.65

# ── Left: Domain-Agnostic components — one box per item ─────────────────────
LX, LW = 0.35, 5.30
section_label(s1, "DOMAIN-AGNOSTIC  —  reusable for any dataset", LX, CONTENT_TOP, LW)

agnostic_items = [
    ("Agent + Skill",
     "Claude Code runs the segment-discovery skill — orchestrates stages, narrates "
     "results, confirms spec with analyst. Path A (guided) and Path B (cold-start) "
     "both live here."),
    ("Tool 1: Single-feature decision tree cuts",
     "tree_cuts.py — fits depth-2 trees on every feature independently; surfaces "
     "BAU and best single-feature thresholds. Always the first stage."),
    ("Tool 2: Multi-feature subgroup search",
     "subgroup_search.py — pysubgroup beam search across thousands of multi-"
     "condition rules; ranks by WRAcc (size × deviation from BAU). The core value-add."),
    ("Tool 3: Stability validation",
     "stability.py — re-checks each top rule within every time/cohort window; "
     "flags unstable segments. No-op when no time column is set."),
    ("Tool 4: SHAP driver analysis",
     "drivers.py — gradient-boosted model + SHAP; confirms rules are "
     "mechanistically grounded, not data artifacts."),
    ("Orchestrator + PPT builder",
     "run_demo.py chains all stages → outputs/REPORT.md. build_ppt.py generates "
     "the branded deck. Both are fully data-agnostic."),
    ("JSON run spec contract",
     "spec.template.json defines what to provide (target, features, sentinels…) "
     "but never which values. Swap the spec, the rest is unchanged."),
]

item_top = 1.06
GAP = 0.04
item_h = (CONTENT_BOT - item_top - GAP * (len(agnostic_items) - 1)) / len(agnostic_items)
for title, body in agnostic_items:
    gold_bar_card(s1, LX, item_top, LW, item_h)
    txt(s1, title, LX + 0.16, item_top + 0.03, LW - 0.28, 0.16, size=8, bold=True, color=NAVY)
    txt(s1, body,  LX + 0.16, item_top + 0.20, LW - 0.28, item_h - 0.23, size=6.7, color=NAVY)
    item_top += item_h + GAP

# ── Right: Example-Specific inputs — exactly 2 boxes ────────────────────────
RX, RW = 5.95, 3.70
section_label(s1, "EXAMPLE-SPECIFIC  —  provided by the analyst", RX, CONTENT_TOP, RW)

card_h = (CONTENT_BOT - 1.06 - 0.12) / 2
card_top = 1.06

gold_bar_card(s1, RX, card_top, RW, card_h)
txt(s1, "Config File  —  config.json", RX + 0.2, card_top + 0.08, RW - 0.35, 0.24,
    size=10.5, bold=True, color=NAVY)
txt(s1, "structured & technical", RX + 0.2, card_top + 0.30, RW - 0.35, 0.16,
    size=7, italic=True, color=GRAY)
bullets(s1, [
    "Data file + dictionary (any CSV/Excel — tools don't care about column names)",
    "Target derivation rule (e.g. responded = 1 where an ID is not null)",
    "Feature set (e.g. 26 EFX attributes) + missing/sentinel codes",
    "Leakage exclusions + time/stability column for validation",
], RX + 0.2, card_top + 0.5, RW - 0.35, card_h - 0.55, size=7.3)

card2_top = card_top + card_h + 0.12
gold_bar_card(s1, RX, card2_top, RW, card_h)
txt(s1, "Prompt  —  instructions to Claude", RX + 0.2, card2_top + 0.08, RW - 0.35, 0.24,
    size=10.5, bold=True, color=NAVY)
txt(s1, "business & conversational", RX + 0.2, card2_top + 0.30, RW - 0.35, 0.16,
    size=7, italic=True, color=GRAY)
bullets(s1, [
    "Business goal & direction (suppress vs target)",
    "Domain context the analyst already knows",
    "What “good” looks like / review criteria",
    "Ad-hoc guidance not worth encoding in JSON",
], RX + 0.2, card2_top + 0.5, RW - 0.35, card_h - 0.55, size=7.3)

callout(s1, "Today: analyst supplies both manually (Path A)   →   "
            "Next: Path B derives the config file itself from a one-line prompt (Slide 2)")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PATH B: NEXT STEPS
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(LAYOUT)
bg(s2)
page_header(s2,
    "Next Steps: Path B — Cold-Start Without Manual Config",
    "Six capabilities so the agent builds the config file itself when the analyst gives only a prompt")
footer(s2, 2)

steps = [
    (1, "Auto-Inspect",
     "Scans dtypes, null rates, cardinality, ID/date detection.",
     "manual column review"),
    (2, "Dictionary Parser",
     "Extracts sentinel codes & attribute categories from any dictionary file.",
     "hand-entering sentinel codes"),
    (3, "Target + Leakage Detector",
     "Proposes target candidates and flags leakage columns to exclude.",
     "manual target / exclude selection"),
    (4, "Domain Research Agent",
     "Web-searches the problem type for additional useful features.",
     "analyst's domain knowledge"),
    (5, "Auto Spec Writer",
     "Assembles steps 1–4 into config.json for one-step approval.",
     "manual config authoring"),
    (6, "Multi-Agent + Memory",
     "Parallel sub-agents + persisted segment rules across runs.",
     "sequential re-discovery each campaign"),
]

GRID_TOP, GRID_BOT = 0.80, 4.65
ROW_GAP, COL_GAP = 0.15, 0.15
CARD_W = (9.30 - 2 * COL_GAP) / 3
CARD_H = (GRID_BOT - GRID_TOP - ROW_GAP) / 2

for idx, (num, title, desc, replaces) in enumerate(steps):
    row, col = divmod(idx, 3)
    lx = 0.35 + col * (CARD_W + COL_GAP)
    ty = GRID_TOP + row * (CARD_H + ROW_GAP)

    box(s2, lx, ty, CARD_W, CARD_H, WHITE)
    hdr_h = 0.40
    box(s2, lx, ty, CARD_W, hdr_h, NAVY)
    box(s2, lx + 0.08, ty + 0.06, 0.28, 0.28, GOLD)
    txt(s2, str(num), lx + 0.08, ty + 0.055, 0.28, 0.28, size=10, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER)
    txt(s2, title, lx + 0.44, ty + 0.06, CARD_W - 0.54, 0.3, size=9.5, bold=True, color=WHITE)

    txt(s2, desc, lx + 0.12, ty + hdr_h + 0.08, CARD_W - 0.24, CARD_H - hdr_h - 0.5,
        size=7.5, color=NAVY)

    div_y = ty + CARD_H - 0.34
    box(s2, lx + 0.12, div_y, CARD_W - 0.24, 0.014, GOLD)
    txt(s2, f"Replaces: {replaces}", lx + 0.12, div_y + 0.04, CARD_W - 0.24, 0.28,
        size=6.8, italic=True, color=GRAY)

callout(s2, "End state: prompt only → agent inspects, researches, drafts config.json, "
            "confirms, and runs")

prs.save(OUT)
print(f"Saved -> {OUT}")
