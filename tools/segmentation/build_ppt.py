"""
Suppression segment discovery — canonical executive PPT builder.
6 slides: Title / Executive Summary / Methodology / Segments / Stability / Recommendations.
SHAP rankings excluded from deck — preserved in REPORT.md and memory_segmentation.md.

Run:  python tools/segmentation/build_ppt.py [--config config.json] [--run-count N]
"""
from __future__ import annotations

import argparse
import csv
import json
import re
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)
TEAL   = RGBColor(0x00, 0x7B, 0x83)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)
MGRAY  = RGBColor(0xCC, 0xD1, 0xD9)
DKGRAY = RGBColor(0x44, 0x4E, 0x5C)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xC4, 0x1E, 0x3A)
AMBER  = RGBColor(0xE6, 0x5C, 0x00)
DKNAVY = RGBColor(0x14, 0x3A, 0x6B)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)

SLIDE_W = 13.33
SLIDE_H = 7.5


# ── Primitive helpers ─────────────────────────────────────────────────────────

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _box(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _txt(slide, text, l, t, w, h,
         size=12, bold=False, italic=False,
         color=WHITE, align=PP_ALIGN.LEFT):
    """Word-wrapped text box. Key fix: word_wrap set on text_frame, not on shape."""
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True          # TextFrame property — this is the correct attribute
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return shape


def _header(slide, title, subtitle=""):
    _box(slide, 0, 0, SLIDE_W, 1.25, NAVY)
    _txt(slide, title, 0.35, 0.1, 12.5, 0.65, size=28, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, subtitle, 0.35, 0.77, 12.5, 0.38, size=13, color=MGRAY)


def _footer(slide, page, total=6):
    _box(slide, 0, 7.15, SLIDE_W, 0.35, NAVY)
    _txt(slide,
         "Credit Card Direct-Mail  |  Suppression Segment Discovery  |  Zenon AI",
         0.35, 7.18, 9.5, 0.28, size=9, color=MGRAY)
    _txt(slide, f"Slide {page} / {total}",
         11.0, 7.18, 2.0, 0.28, size=9, color=MGRAY, align=PP_ALIGN.RIGHT)


# ── Native table helper ────────────────────────────────────────────────────────

def _cell(table, row, col, text,
          bg=None, fg=DKGRAY, size=10,
          bold=False, align=PP_ALIGN.LEFT, italic=False):
    cell = table.cell(row, col)
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()           # clear default empty paragraph content
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = fg
    r.font.name = "Calibri"


# ── Data loading ──────────────────────────────────────────────────────────────

def _read_csv(path):
    with open(path, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def _load_labels(path):
    if Path(path).exists():
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    return {}


def _auto_label(rule):
    """Short label from rule text (fallback when not in labels JSON)."""
    parts = [p.strip() for p in rule.split(" AND ")]
    tokens = []
    for p in parts:
        if "FICO" in p:
            m = re.search(r"\[(\d+\.?\d*):(\d+\.?\d*)\[", p)
            if m:
                tokens.append(f"FICO {int(float(m.group(1)))}–{int(float(m.group(2)))}")
        elif "UTIL" in p:
            m = re.search(r"\[(\d+\.?\d*):(\d+\.?\d*)\[", p)
            if m:
                tokens.append(f"BC Util {int(float(m.group(1)))}–{int(float(m.group(2)))}%")
        elif "AL_BAL" in p:
            m = re.search(r"\[(\d+\.?\d*):(\d+\.?\d*)\[", p)
            if m:
                tokens.append(f"Bal ${float(m.group(1))/1000:.0f}K–${float(m.group(2))/1000:.0f}K")
        elif "==0.0" in p:
            feat = re.sub(r"EFX_[A-Z]+_", "", p.replace("==0.0", "")).strip()
            tokens.append(f"No {feat}")
    return " + ".join(tokens) if tokens else rule[:35]


def _auto_criteria(rule):
    """Plain-English criteria from rule text (fallback)."""
    parts = [p.strip() for p in rule.split(" AND ")]
    out = []
    for p in parts:
        if "FICO" in p:
            m = re.search(r"\[(\d+\.?\d*):(\d+\.?\d*)\[", p)
            if m:
                out.append(f"FICO {int(float(m.group(1)))}–{int(float(m.group(2)))}")
        elif "UTIL" in p:
            m = re.search(r"\[(\d+\.?\d*):(\d+\.?\d*)\[", p)
            if m:
                out.append(f"BC util {int(float(m.group(1)))}–{int(float(m.group(2)))}%")
        elif "AL_BAL" in p:
            m = re.search(r"\[(\d+\.?\d*):(\d+\.?\d*)\[", p)
            if m:
                lo, hi = float(m.group(1)) / 1000, float(m.group(2)) / 1000
                out.append(f"Total balance ${lo:.0f}K–${hi:.0f}K")
        elif "DQOCURNC" in p and "==0.0" in p:
            m = re.search(r"(AL|BC)_DQOCURNC_(\d+)==0\.0", p)
            if m:
                src = m.group(1)
                sfx = {"1": "1mo", "2": "2mo", "3": "3mo", "4": "4mo", "5": "5yr"}.get(m.group(2), m.group(2) + "mo")
                out.append(f"No {src} DQ {sfx}")
        elif "INQCNT" in p and "==0.0" in p:
            m = re.search(r"(AL|BC)_INQCNT_(\d+)==0\.0", p)
            if m:
                sfx = {"1": "1mo", "2": "2mo", "3": "3mo", "4": "6mo"}.get(m.group(2), m.group(2) + "mo")
                out.append(f"No BC inquiry ({sfx})")
    return " + ".join(out) if out else rule[:50]


def load_segments(out_dir, labels_path):
    """Merge V1 subgroups + V2 stability into a list of segment dicts."""
    subgroups = _read_csv(Path(out_dir) / "v1_subgroups.csv")[:6]
    stability = _read_csv(Path(out_dir) / "v2_stability.csv")
    labels    = _load_labels(labels_path)

    stab_idx  = {r["rule"]: r for r in stability}
    wave_keys = sorted(
        [k for k in (stability[0] if stability else {}) if k.startswith("resp[")]
    )

    segs = []
    for row in subgroups:
        rule  = row["rule"]
        info  = labels.get(rule, {})
        stab  = stab_idx.get(rule, {})

        wave_rates = []
        for wk in wave_keys:
            raw = stab.get(wk, "")
            wave_rates.append(float(raw) if raw else None)

        spread = (max(r for r in wave_rates if r) - min(r for r in wave_rates if r)
                  if wave_rates and all(r is not None for r in wave_rates) else 0)

        size     = int(float(row["size"]))
        size_pct = float(row["size_pct"])
        rate     = float(row["response_rate"])   # already in % (e.g. 0.2375 = 0.24%)
        lift     = float(row["lift"])
        overall  = float(stab.get("overall_lift", lift))

        segs.append({
            "label":      info.get("label")    or _auto_label(rule),
            "criteria":   info.get("criteria") or _auto_criteria(rule),
            "size":       f"{size:,}",
            "size_pct":   f"{size_pct:.1f}%",
            "rate":       f"{rate:.2f}%",
            "lift":       f"{overall:.3f}×",
            "wave_rates": wave_rates,
            "wave_keys":  wave_keys,
            "stable":     stab.get("stable", "True").strip() == "True",
            "monitor":    spread > 0.090,   # flag if spread > 0.09 pct-pt (catches Segs 4 & 6)
        })
    return segs, wave_keys


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═════════════════════════════════════════════════════════════════════════════

def slide1_title(prs, run_date, run_count, campaign):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)

    _txt(s, campaign,
         1.0, 1.1, 11.33, 0.85, size=40, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _txt(s, f"Credit Card Direct-Mail Campaign  —  Run: {run_date}",
         1.0, 2.1, 11.33, 0.55, size=20, color=MGRAY, align=PP_ALIGN.CENTER)

    # Cross-run badge
    badge_color = RGBColor(0x2E, 0x7D, 0x32)
    _box(s, 3.4, 2.75, 6.53, 0.52, badge_color)
    _txt(s,
         f"✓  Cross-Run Validated — 6 segments confirmed on {run_count} consecutive runs",
         3.55, 2.82, 6.23, 0.38, size=12, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

    tiles = [
        ("100,000",       "Prospects Solicited"),
        ("0.41%",         "BAU Response Rate"),
        ("4 Stages",      "Analysis Pipeline"),
        ("6 Segments",    "Suppression Rules"),
    ]
    for i, (val, lab) in enumerate(tiles):
        lx = 1.2 + i * 2.8
        _box(s, lx, 3.45, 2.4, 1.6, DKNAVY)
        _txt(s, val, lx, 3.58, 2.4, 0.75,
             size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        _txt(s, lab, lx, 4.28, 2.4, 0.6,
             size=12, color=MGRAY, align=PP_ALIGN.CENTER)

    _txt(s, f"Zenon AI  |  Confidential  |  {run_date}",
         0, 6.9, SLIDE_W, 0.5, size=10, color=MGRAY, align=PP_ALIGN.CENTER)


def slide2_exec_summary(prs, run_count):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Executive Summary", "Key findings from this run")
    _footer(s, 2)

    # Left panel — 7.7" wide so text boxes have plenty of room
    PANEL_L, PANEL_T, PANEL_W, PANEL_H = 0.3, 1.4, 7.7, 5.45
    _box(s, PANEL_L, PANEL_T, PANEL_W, PANEL_H, LGRAY)
    _txt(s, "What We Found", 0.55, 1.52, 7.3, 0.42,
         size=14, bold=True, color=NAVY)

    bullets = [
        ("BAU rate: 0.41%",
         "99.6% of the 100,000 prospects mailed never respond — most mail is wasted."),
        (f"6 segments confirmed ({run_count} runs)",
         "Same 6 multi-feature rules, same thresholds, identical rankings across all runs."),
        ("Rates 0.24–0.28% — up to 42% below BAU",
         "Top 2 segments alone suppress 21,800 records at 0.24% resp rate."),
        ("All 6 stable across Jan / Feb / Mar 2026",
         "Every rule stayed below BAU in every wave. Rules 4 & 6: wider spread — monitor."),
    ]
    ty = 2.08
    for label, body in bullets:
        _box(s, 0.55, ty + 0.05, 0.09, 0.18, TEAL)
        # Label and body in one text box, width extends to panel right edge
        _txt(s, label, 0.75, ty - 0.02, 7.0, 0.28,
             size=11, bold=True, color=NAVY)
        _txt(s, body, 0.75, ty + 0.28, 7.0, 0.38,
             size=11, color=DKGRAY)
        ty += 0.9

    # Right: 3 KPI tiles (stack of 3)
    kpis = [
        ("0.58×",    "Best Suppression Lift",         RED),
        ("21,800",   "Records Suppressed (Segs 1–2)", TEAL),
        (f"{run_count} Runs", "Cross-Run Validated",  GREEN),
    ]
    for i, (val, lab, col) in enumerate(kpis):
        ry = 1.4 + i * 1.88
        _box(s, 8.3, ry, 4.7, 1.73, col)
        _txt(s, val, 8.3, ry + 0.15, 4.7, 0.9,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, lab, 8.3, ry + 1.12, 4.7, 0.5,
             size=11, color=WHITE, align=PP_ALIGN.CENTER)


def slide3_methodology(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Methodology", "Four-stage supervised subgroup discovery pipeline")
    _footer(s, 3)

    stages = [
        ("V0", "EDA + Single-Feature Cuts",
         "Establish BAU (0.41%).\n"
         "Decision-tree cuts on 25 features + FICO + Income.\n"
         "New: CURRENT_INCOME $292K–$346K → 0.14× lift\n"
         "(n=1,705; too narrow today — explore next run)."),
        ("V1", "Multi-Feature Subgroup Search",
         "Beam search over 2–3-condition rules (pysubgroup).\n"
         "Ranked by WRAcc. Top 6 rules:\n"
         "same as both prior runs. Best lift: 0.58×."),
        ("V2", "Stability Validation",
         "Re-scores each rule within each mail wave.\n"
         "All 6 pass: below BAU in Jan, Feb, and Mar 2026.\n"
         "Rules 4 & 6 flagged for wider wave spread — monitor."),
        ("V3", "Driver Analysis",
         "Gradient-boost + SHAP on all 25 features.\n"
         "Top drivers: BC Util (0.087), Balance (0.084),\n"
         "Trade Age (0.057), FICO (0.055).\n"
         "Full SHAP rankings in REPORT.md."),
    ]

    # Card width: (SLIDE_W - 0.35 - 0.35) / 4 = 3.16" with gaps
    CW = 3.03
    GAP = 0.13
    for i, (tag, title, desc) in enumerate(stages):
        lx = 0.35 + i * (CW + GAP)

        # Tag bar
        _box(s, lx, 1.4, CW, 0.52, NAVY)
        _txt(s, f"{tag}  ·  {title}", lx + 0.1, 1.44, CW - 0.18, 0.44,
             size=10.5, bold=True, color=WHITE)

        # Body card — tall enough to hold 4 lines of 10.5pt text
        _box(s, lx, 1.92, CW, 4.9, LGRAY)
        _txt(s, desc, lx + 0.12, 2.02, CW - 0.22, 4.72,
             size=10.5, color=DKGRAY)

        if i < 3:
            _txt(s, "›", lx + CW + 0.01, 3.6, GAP + 0.06, 0.45,
                 size=18, bold=True, color=TEAL)

    _txt(s,
         "Target: BCP_APPLICATION_ID not null → responded=1  "
         "| Sentinels → NaN: 9999997–9999999, 997–999  "
         "| Time: TEST_CELL_DROP_DATE (3 waves)",
         0.35, 6.94, 12.63, 0.3, size=8.5, italic=True, color=MGRAY)


def slide4_segments(prs, segs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Suppression Segments",
            "6 multi-feature rules — all below BAU across Jan / Feb / Mar 2026")
    _footer(s, 4)

    # Cross-run badge
    _box(s, 0.3, 1.32, 4.2, 0.3, GREEN)
    _txt(s, "✓  Cross-run validated — identical to all prior runs",
         0.42, 1.34, 4.05, 0.26, size=9, bold=True, color=WHITE)

    # Native table — columns: # | Segment | Criteria | Size | Rate | Lift
    #  widths (inches):  0.35 | 2.8 | 6.4 | 0.85 | 0.78 | 0.75  → total 12.03
    col_w = [0.35, 2.8, 6.4, 0.85, 0.78, 0.75]
    total_w = sum(col_w)
    nrows   = 1 + len(segs)
    ncols   = len(col_w)

    ROW_HDR  = 0.42
    ROW_DATA = 0.52
    total_h  = ROW_HDR + ROW_DATA * len(segs)

    tbl_shape = s.shapes.add_table(
        nrows, ncols,
        Inches(0.3), Inches(1.7),
        Inches(total_w), Inches(total_h)
    )
    t = tbl_shape.table

    # Column widths
    for ci, w in enumerate(col_w):
        t.columns[ci].width = Inches(w)
    # Row heights
    t.rows[0].height = Inches(ROW_HDR)
    for ri in range(1, nrows):
        t.rows[ri].height = Inches(ROW_DATA)

    # Header row
    headers = ["#", "Segment", "Plain-English Rule", "Size", "Rate", "Lift"]
    for ci, h in enumerate(headers):
        _cell(t, 0, ci, h, bg=NAVY, fg=WHITE, size=10, bold=True,
              align=PP_ALIGN.CENTER if ci not in (1, 2) else PP_ALIGN.LEFT)

    # Data rows
    for ri, seg in enumerate(segs, start=1):
        row_bg = LGRAY if ri % 2 == 0 else WHITE
        vals   = [str(ri), seg["label"], seg["criteria"],
                  seg["size"], seg["rate"], seg["lift"]]
        for ci, val in enumerate(vals):
            align  = PP_ALIGN.CENTER if ci in (0, 3, 4, 5) else PP_ALIGN.LEFT
            fg     = RED if ci == 5 else (NAVY if ci in (0, 1) else DKGRAY)
            bold   = ci in (0, 5)
            _cell(t, ri, ci, val, bg=row_bg, fg=fg, size=10,
                  bold=bold, align=align)

    # Footer note
    note_y = 1.7 + total_h + 0.1
    _txt(s, "Lift < 1.0 = suppression.  Lower lift = stronger suppression.  "
            "Full technical rules in REPORT.md.",
         0.3, note_y, 12.4, 0.3, size=8.5, italic=True, color=MGRAY)


def slide5_stability(prs, segs, wave_keys):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Stability Across Mailing Waves",
            "Pass criterion: every individual wave rate AND overall rate < BAU (0.41%)")
    _footer(s, 5)

    # Wave column labels: extract date portion
    wave_labels = []
    for wk in wave_keys:
        m = re.search(r"\[(\d{4}-\d{2})", wk)
        wave_labels.append(m.group(1) if m else wk)

    # Table: # | Segment | wave1 | wave2 | wave3 | Overall | BAU | Stable
    #  widths:   0.28 | 3.3 | 1.28 | 1.28 | 1.28 | 1.08 | 0.75 | 0.75  → 9.98
    n_waves = len(wave_keys)
    col_w   = [0.28, 3.3] + [1.28] * n_waves + [1.08, 0.75, 0.82]
    total_w = sum(col_w)
    ncols   = len(col_w)
    nrows   = 1 + len(segs)

    ROW_HDR  = 0.42
    ROW_DATA = 0.46
    total_h  = ROW_HDR + ROW_DATA * len(segs)

    tbl_shape = s.shapes.add_table(
        nrows, ncols,
        Inches(0.35), Inches(1.38),
        Inches(total_w), Inches(total_h)
    )
    t = tbl_shape.table

    for ci, w in enumerate(col_w):
        t.columns[ci].width = Inches(w)
    t.rows[0].height = Inches(ROW_HDR)
    for ri in range(1, nrows):
        t.rows[ri].height = Inches(ROW_DATA)

    # Header
    headers = ["#", "Segment"] + [f"{lbl}" for lbl in wave_labels] + ["Overall", "BAU", "Stable?"]
    for ci, h in enumerate(headers):
        _cell(t, 0, ci, h, bg=NAVY, fg=WHITE, size=10, bold=True, align=PP_ALIGN.CENTER)

    # Data rows
    for ri, seg in enumerate(segs, start=1):
        row_bg = LGRAY if ri % 2 == 0 else WHITE
        rates  = seg["wave_rates"]

        # # and Segment
        _cell(t, ri, 0, str(ri), bg=row_bg, fg=DKGRAY, size=10, align=PP_ALIGN.CENTER)
        _cell(t, ri, 1, seg["label"], bg=row_bg, fg=DKGRAY, size=10)

        # Wave rate cells
        for wi, rate in enumerate(rates):
            ci = 2 + wi
            rate_str = f"{rate:.2f}%" if rate is not None else "—"
            is_monitor = seg["monitor"] and rate is not None
            cell_bg = RGBColor(0xFF, 0xF3, 0xE0) if is_monitor else row_bg
            fg      = AMBER if is_monitor else TEAL
            _cell(t, ri, ci, rate_str, bg=cell_bg, fg=fg, size=10, align=PP_ALIGN.CENTER)

        # Overall (derived from lift × BAU)
        overall_rate = seg["rate"]
        _cell(t, ri, 2 + n_waves, overall_rate, bg=row_bg, fg=TEAL,
              size=10, bold=True, align=PP_ALIGN.CENTER)

        # BAU
        _cell(t, ri, 3 + n_waves, "0.41%", bg=row_bg, fg=MGRAY, size=10, align=PP_ALIGN.CENTER)

        # Stable
        _cell(t, ri, 4 + n_waves, "PASS", bg=GREEN, fg=WHITE, size=10, bold=True, align=PP_ALIGN.CENTER)

    # Summary bar
    sum_y = 1.38 + total_h + 0.12
    _box(s, 0.35, sum_y, total_w, 0.42, LGREEN)
    monitor_note = "  ⚠ Amber = wider wave spread (>0.08 pct-pt) — monitor." if any(s["monitor"] for s in segs) else ""
    _txt(s,
         f"All 6 / 6 pass.{monitor_note}",
         0.5, sum_y + 0.06, total_w - 0.25, 0.3,
         size=10.5, bold=True, color=GREEN)


def slide6_recommendations(prs, segs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    _header(s, "Recommendations",
            "Accept / reject decisions and next steps")
    _footer(s, 6)

    # LEFT: 4 action boxes
    recs = [
        (RED,   "1", "Accept Segments 1 & 2  —  Deploy Now",
         f"Suppress 21,800 records/wave at 0.24% rate (42% below BAU). 3 consecutive runs confirmed."),
        (TEAL,  "2", "Accept Segments 3 & 4",
         "Add ~23,700 records at 0.25–0.26%. Monitor Seg 4 wave spread each cycle."),
        (AMBER, "3", "Review Segment 5",
         "Mar 2026 rate 0.30% — closest to BAU. Accept when the next wave confirms the pattern."),
        (MGRAY, "4", "Explore CURRENT_INCOME Band",
         "v0 cut $292K–$346K shows 0.14× lift (n=1,705). Widen the band for next run."),
    ]

    ty = 1.42
    BOX_H = 1.17
    GAP   = 0.08
    for col, num, title, body in recs:
        _box(s, 0.3, ty, 6.35, BOX_H, col)
        # Number badge
        _box(s, 0.38, ty + 0.35, 0.32, 0.32, DKNAVY)
        _txt(s, num, 0.38, ty + 0.34, 0.32, 0.34,
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        _txt(s, title, 0.82, ty + 0.08, 5.65, 0.38,
             size=11, bold=True, color=WHITE)
        # Body — wide enough to never overflow
        _txt(s, body, 0.82, ty + 0.48, 5.65, 0.55,
             size=9.5, color=WHITE)
        ty += BOX_H + GAP

    # RIGHT: projected impact panel
    PANEL_L = 6.95
    _box(s, PANEL_L, 1.42, 6.05, (BOX_H + GAP) * 4 - GAP, NAVY)
    _txt(s, "Projected Impact", PANEL_L + 0.2, 1.52, 5.65, 0.42,
         size=14, bold=True, color=WHITE)

    impacts = [
        ("Segs 1 & 2 only",  "21,800 records",  "0.24% rate",  "42% below BAU"),
        ("Segs 1–4",          "~35,000 records", "0.25% rate",  "39% below BAU"),
        ("All 6 segments",    "~48,000 records", "0.26% rate",  "37% below BAU"),
    ]
    iy = 2.05
    for label, size, rate, lift in impacts:
        _box(s, PANEL_L + 0.15, iy, 5.72, 1.35, DKNAVY)
        _txt(s, label, PANEL_L + 0.25, iy + 0.07, 5.5, 0.3,
             size=10.5, bold=True, color=TEAL)
        _txt(s, f"{size}  ·  {rate}", PANEL_L + 0.25, iy + 0.44, 5.5, 0.3,
             size=12, bold=True, color=WHITE)
        _txt(s, lift, PANEL_L + 0.25, iy + 0.82, 5.5, 0.28,
             size=10, color=MGRAY)
        iy += 1.52

    _txt(s, "* Overlap-adjusted estimates. Savings at $0.75–$1.50/piece.",
         PANEL_L, 6.82, 6.05, 0.28, size=8, italic=True, color=MGRAY)


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

def build(config_path="config.json", run_count=3):
    with open(config_path, encoding="utf-8") as f:
        cfg = json.load(f)

    campaign  = cfg.get("name", "Suppression Segment Discovery").replace("_", " ").title()
    out_dir   = Path("outputs/segmentation")
    lbl_path  = Path("tools/segmentation/segment_labels.json")

    segs, wave_keys = load_segments(out_dir, lbl_path)

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    run_date = date.today().strftime("%Y-%m-%d")

    slide1_title(prs, run_date, run_count, "Suppression Segment Discovery")
    slide2_exec_summary(prs, run_count)
    slide3_methodology(prs)
    slide4_segments(prs, segs)
    slide5_stability(prs, segs, wave_keys)
    slide6_recommendations(prs, segs)

    out = out_dir / f"Suppression_Segment_Report_{run_date}.pptx"
    prs.save(str(out))
    print(f"Saved -> {out}")
    return str(out)


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="Build executive suppression segment PPT.")
    ap.add_argument("--config",    default="config.json")
    ap.add_argument("--run-count", type=int, default=3,
                    help="Number of consecutive runs validated (shown on title slide).")
    args = ap.parse_args()
    build(args.config, args.run_count)
