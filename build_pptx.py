"""Build executive-ready suppression segment discovery PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Brand palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy  – headings / header bars
BLUE   = RGBColor(0x20, 0x5C, 0xA8)   # mid blue   – accents
RED    = RGBColor(0xC0, 0x39, 0x2B)   # alert red  – suppression highlight
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)   # green      – positive metrics
LGRAY  = RGBColor(0xF4, 0xF6, 0xF8)   # light grey – slide bg panels
DGRAY  = RGBColor(0x5A, 0x6A, 0x7A)   # dark grey  – body text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xE6, 0x9B, 0x00)

TODAY  = datetime.date.today().strftime("%B %d, %Y")

W, H   = Inches(13.33), Inches(7.5)   # widescreen 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────────────
def rgb(r):   return RGBColor(*[int(r[i:i+2], 16) for i in (0, 2, 4)])

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.line.fill.background()
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.0, fill=NAVY)
    add_text(slide, title, 0.3, 0.08, 11, 0.55, size=24, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.3, 0.62, 11, 0.32, size=12,
                 color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)
    add_rect(slide, 0, 1.0, 13.33, 0.04, fill=BLUE)

def footer(slide, page):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=LGRAY)
    add_text(slide, f"CONFIDENTIAL  |  Credit Card DM Suppression  |  {TODAY}  |  Page {page}",
             0.3, 7.22, 12.5, 0.26, size=8, color=DGRAY, align=PP_ALIGN.LEFT)

def kpi_box(slide, l, t, w, h, label, value, value_color=NAVY, bg=LGRAY):
    add_rect(slide, l, t, w, h, fill=bg, line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
    add_text(slide, label, l+0.1, t+0.08, w-0.2, 0.35, size=10, color=DGRAY,
             align=PP_ALIGN.CENTER)
    add_text(slide, value, l+0.05, t+0.38, w-0.1, 0.65, size=20, bold=True,
             color=value_color, align=PP_ALIGN.CENTER)

def table_row(slide, cols, l, t, row_h, widths, bold=False,
              bg=WHITE, text_color=DGRAY, sizes=None):
    x = l
    for i, (cell, w) in enumerate(zip(cols, widths)):
        add_rect(slide, x, t, w, row_h, fill=bg,
                 line=RGBColor(0xDD, 0xDD, 0xDD), line_w=Pt(0.5))
        fs = sizes[i] if sizes else 9
        add_text(slide, str(cell), x+0.06, t+0.04, w-0.12, row_h-0.08,
                 size=fs, bold=bold, color=text_color, align=PP_ALIGN.CENTER)
        x += w


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 1 – Cover
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(sl, 0, 5.8, 13.33, 0.08, fill=BLUE)
add_rect(sl, 0, 5.88, 13.33, 1.62, fill=RGBColor(0x12, 0x1C, 0x35))

# Title
add_text(sl, "Credit Card Direct Mail", 0.7, 1.2, 12, 1.0,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Suppression Segment Discovery", 0.7, 2.1, 12, 0.9,
         size=32, bold=False, color=RGBColor(0xAA, 0xCC, 0xFF),
         align=PP_ALIGN.LEFT)

add_rect(sl, 0.7, 3.15, 4.5, 0.06, fill=GOLD)
add_text(sl, "Identifying non-responder cohorts to remove from future mail waves",
         0.7, 3.35, 11, 0.5, size=14, color=RGBColor(0xCC, 0xDD, 0xFF),
         align=PP_ALIGN.LEFT, italic=True)

add_text(sl, f"Prepared: {TODAY}", 0.7, 6.05, 6, 0.35, size=11,
         color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.LEFT)
add_text(sl, "CONFIDENTIAL", 9, 6.05, 4, 0.35, size=11,
         color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.RIGHT)

add_text(sl, "Powered by Supervised Subgroup Discovery  ·  Tri-bureau credit attributes  ·  100,000 solicited prospects",
         0.7, 6.7, 12, 0.4, size=9, color=RGBColor(0x66, 0x88, 0xAA),
         align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 2 – Executive Summary
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Executive Summary", "Key findings across 100,000 solicited prospects · 3-wave direct mail campaign")
footer(sl, 2)

# Three KPI boxes
kpi_box(sl, 0.3,  1.25, 2.8, 1.3, "Campaign Base Rate", "0.41%",  value_color=DGRAY)
kpi_box(sl, 3.35, 1.25, 2.8, 1.3, "Top Segment Lift",   "0.58×",  value_color=RED)
kpi_box(sl, 6.4,  1.25, 2.8, 1.3, "Mail Savings (Top 5)", "33.8%", value_color=GREEN)
kpi_box(sl, 9.45, 1.25, 3.55, 1.3, "Responders Retained", "78.3%", value_color=BLUE)

add_rect(sl, 0.3, 2.85, 12.73, 0.04, fill=BLUE)

# Finding bullets
bullets = [
    ("🎯 Suppression Opportunity",
     "5 validated segments totalling 33,777 prospects (33.8% of the mailing file) "
     "respond consistently below the campaign BAU in all three mail waves (Jan, Feb, Mar 2026). "
     "Suppressing these records saves one-third of mailing costs while retaining 78.3% of all responses."),
    ("📊 What Drives Non-Response",
     "SHAP analysis confirms the top 5 predictors are: bankcard utilization, total open balance, "
     "trade age, FICO score, and delinquency recurrence. The suppression profile is a 'passively "
     "creditworthy' prospect: mid-FICO (664–730), no recent inquiries, modest utilization — already "
     "credit-satisfied and unlikely to act on a solicitation."),
    ("✅ Stability Validated",
     "All 5 segments hold below BAU across every drop wave. Response rates range 0.23%–0.26% "
     "vs a 0.41% population baseline, representing a 37–43% relative reduction. "
     "Rules are simple (2–3 bureau conditions) and immediately operationalisable in pre-screen or "
     "address-file logic."),
]
y = 3.0
for title, body in bullets:
    add_rect(sl, 0.3, y, 12.73, 0.3, fill=NAVY)
    add_text(sl, title, 0.35, y+0.02, 12, 0.26, size=11, bold=True, color=WHITE)
    add_text(sl, body, 0.4, y+0.32, 12.3, 0.72, size=10, color=DGRAY)
    add_rect(sl, 0.3, y+0.32, 12.73, 0.72, fill=LGRAY,
             line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
    add_text(sl, body, 0.45, y+0.36, 12.2, 0.65, size=10, color=DGRAY)
    y += 1.12


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 3 – Campaign Baseline & Methodology
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Campaign Baseline & Methodology",
           "100,000 prospects · Tri-bureau (EFX/EXP/TRU) credit attributes · 4-stage analytical pipeline")
footer(sl, 3)

# Left panel – baseline stats
add_rect(sl, 0.3, 1.2, 5.5, 5.85, fill=LGRAY, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 0.3, 1.2, 5.5, 0.45, fill=NAVY)
add_text(sl, "CAMPAIGN BASELINE", 0.35, 1.22, 5, 0.4, size=12, bold=True, color=WHITE)

stats = [
    ("Population", "100,000 solicited prospects"),
    ("Target variable", "responded = BCP_APPLICATION_ID not null"),
    ("Responders", "410 prospects (0.41% BAU)"),
    ("Bookers", "84 prospects (0.084% booking rate)"),
    ("Mail waves", "3 drop dates (Jan 6, Feb 10, Mar 10, 2026)"),
    ("Bureau coverage", "Equifax · Experian · TransUnion (26 EFX features used)"),
    ("Sentinel codes", "Treated as missing: 9999997/98/99, 997/98/99"),
    ("Direction", "SUPPRESSION — find low-response cohorts"),
]
y = 1.75
for k, v in stats:
    add_text(sl, k + ":", 0.45, y, 2.0, 0.3, size=9, bold=True, color=NAVY)
    add_text(sl, v, 2.5, y, 3.2, 0.3, size=9, color=DGRAY)
    y += 0.42

# Right panel – methodology pipeline
add_rect(sl, 6.2, 1.2, 6.83, 5.85, fill=WHITE,
         line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 6.2, 1.2, 6.83, 0.45, fill=BLUE)
add_text(sl, "4-STAGE ANALYTICAL PIPELINE", 6.25, 1.22, 6.5, 0.4,
         size=12, bold=True, color=WHITE)

stages = [
    ("V0", "One-Feature Tree Cuts",
     "Decision-tree splits on each variable independently. "
     "Establishes single-feature suppression thresholds; reveals "
     "which attributes show any signal at all."),
    ("V1", "Multi-Feature Subgroup Search",
     "Beam search over all 2–3 condition rule combinations "
     "(pysubgroup). Finds segments the manual loop misses. "
     "Ranked by quality × size × suppression depth."),
    ("V2", "Stability Validation",
     "Each top rule re-evaluated within each mail wave. "
     "Confirms segments are directionally below BAU across "
     "all three drop dates — not a one-wave artefact."),
    ("V3", "SHAP Driver Analysis",
     "XGBoost + TreeSHAP over all features. Identifies the "
     "global levers that explain non-response. Validates that "
     "segment conditions map to mechanistically sensible drivers."),
]
y = 1.75
for code, title, desc in stages:
    add_rect(sl, 6.3, y, 0.55, 0.9, fill=BLUE)
    add_text(sl, code, 6.3, y+0.2, 0.55, 0.5, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 6.95, y+0.02, 5.9, 0.3, size=10, bold=True, color=NAVY)
    add_text(sl, desc, 6.95, y+0.32, 5.95, 0.6, size=9, color=DGRAY)
    if y < 4.5:
        add_rect(sl, 6.3, y+0.92, 0.01, 0.3, fill=BLUE)
    y += 1.3


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 4 – Top 5 Suppression Segments
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Top 5 Validated Suppression Segments",
           "All 5 segments hold below BAU in every mail wave · V1 multi-feature subgroup search · ranked by suppression depth")
footer(sl, 4)

# Table header
cols_h = ["#", "Segment Rule (plain English)", "Size", "% of File",
          "Resp Rate", "BAU", "Lift", "Status"]
widths  = [0.4, 5.5, 0.85, 0.75, 0.85, 0.65, 0.75, 1.08]
table_row(sl, cols_h, 0.3, 1.2, 0.38, widths, bold=True,
          bg=NAVY, text_color=WHITE, sizes=[9,9,9,9,9,9,9,9])

segments = [
    ("1", "No long-term DQ (all trades) +\nNo bankcard DQ mid-term +\nFICO 664–694",
     "10,947", "10.95%", "0.24%", "0.41%", "0.58×"),
    ("2", "No recent DQ (all trades) +\nNo bankcard inquiries (6 mo) +\nBC util 16–35%",
     "10,881", "10.88%", "0.24%", "0.41%", "0.58×"),
    ("3", "No bankcard DQ mid-term +\nNo bankcard inquiries (6 mo) +\nFICO 664–694",
     "10,616", "10.62%", "0.25%", "0.41%", "0.62×"),
    ("4", "No bankcard DQ long-term +\nNo bankcard inquiries (6 mo) +\nFICO 694–730",
     "13,060", "13.06%", "0.26%", "0.41%", "0.64×"),
    ("5", "No bankcard DQ mid-term +\nFICO 664–694",
     "12,939", "12.94%", "0.26%", "0.41%", "0.64×"),
]

row_colors = [LGRAY, WHITE, LGRAY, WHITE, LGRAY]
y = 1.58
for i, (num, rule, size, pct, rr, bau_v, lift) in enumerate(segments):
    bg = row_colors[i]
    h = 0.78
    add_rect(sl, 0.3, y, 0.4, h, fill=NAVY)
    add_text(sl, num, 0.3, y+0.22, 0.4, 0.35, size=13, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 0.7, y, 5.5, h, fill=bg,
             line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
    add_text(sl, rule, 0.78, y+0.04, 5.3, 0.72, size=9, color=DGRAY)
    for val, ww in zip([size, pct, rr, bau_v, lift],
                       [0.85, 0.75, 0.85, 0.65, 0.75]):
        x_off = sum(widths[:widths.index(ww)+1]) + 0.3 - ww
        # use offset accumulation
    x = 0.3 + 0.4 + 5.5
    for val, ww in zip([size, pct, rr, bau_v, lift], [0.85, 0.75, 0.85, 0.65, 0.75]):
        add_rect(sl, x, y, ww, h, fill=bg,
                 line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
        vc = RED if val == lift else DGRAY
        add_text(sl, val, x+0.04, y+0.2, ww-0.08, 0.38,
                 size=10, bold=(val==lift), color=vc, align=PP_ALIGN.CENTER)
        x += ww

    # Status badge
    add_rect(sl, x, y, 1.08, h, fill=GREEN)
    add_text(sl, "✓ STABLE\n(all 3 waves)", x+0.05, y+0.12, 0.98, 0.55,
             size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    y += 0.82

# Summary callout
add_rect(sl, 0.3, 5.72, 12.73, 0.72, fill=RGBColor(0xE8, 0xF4, 0xE8),
         line=GREEN, line_w=Pt(1.0))
add_text(sl, "UNION OF TOP 5 SEGMENTS: 33,777 records (33.8% of file) · 89 actual responders vs 138 expected at BAU · "
         "Suppressing these saves ≈34% of mailing cost while retaining 78.3% of all campaign responses.",
         0.45, 5.78, 12.4, 0.62, size=10, bold=False, color=GREEN)


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 5 – Segment 1 Deep Dive
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Segment #1 Deep Dive — Best Suppression Rule",
           "EFX_AL_DQOCURNC_5==0 AND EFX_BC_DQOCURNC_3==0 AND FICO ∈ [664, 694)")
footer(sl, 5)

# KPI strip
for i, (lbl, val, vc) in enumerate([
    ("Segment Size", "10,947", NAVY),
    ("% of File", "10.95%", NAVY),
    ("Response Rate", "0.24%", RED),
    ("vs BAU", "0.41%", DGRAY),
    ("Lift Ratio", "0.58×", RED),
    ("Responders Forfeited", "26 of 410", GOLD),
]):
    kpi_box(sl, 0.3 + i*2.15, 1.25, 2.0, 1.2, lbl, val, value_color=vc)

# Rule translation
add_rect(sl, 0.3, 2.7, 5.9, 3.85, fill=LGRAY, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 0.3, 2.7, 5.9, 0.4, fill=BLUE)
add_text(sl, "RULE IN PLAIN ENGLISH", 0.35, 2.72, 5.5, 0.36, size=11, bold=True, color=WHITE)

conditions = [
    ("Condition 1", "EFX_AL_DQOCURNC_5 = 0",
     "Zero delinquency events across ALL trade lines in the past 36+ months. "
     "This prospect has never missed a payment in the review window."),
    ("Condition 2", "EFX_BC_DQOCURNC_3 = 0",
     "Zero bankcard-specific delinquency events in the 24–36 month window. "
     "Reinforces pattern: no credit stress history."),
    ("Condition 3", "FICO ∈ [664, 694)",
     "Mid near-prime FICO band. Above sub-prime (where cards are scarce), "
     "below prime (where competing offers flood the mailbox). Comfortable "
     "credit access — not motivated by need."),
]
y = 3.18
for code, rule, desc in conditions:
    add_rect(sl, 0.35, y, 1.0, 0.82, fill=BLUE)
    add_text(sl, code.replace("Condition ", "C"), 0.35, y+0.18, 1.0, 0.45,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, rule, 1.42, y+0.02, 4.65, 0.3, size=9, bold=True, color=NAVY)
    add_text(sl, desc, 1.42, y+0.3, 4.65, 0.54, size=8.5, color=DGRAY)
    y += 0.97

# Wave stability table
add_rect(sl, 6.5, 2.7, 6.53, 3.85, fill=WHITE, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 6.5, 2.7, 6.53, 0.4, fill=NAVY)
add_text(sl, "WAVE-BY-WAVE STABILITY", 6.55, 2.72, 6.3, 0.36, size=11, bold=True, color=WHITE)

wave_data = [
    ("Wave", "Drop Date", "Segment N", "Resp Rate", "vs BAU", "Lift"),
    ("W1", "Jan 6, 2026",  "3,650", "0.247%", "0.41%", "0.60×"),
    ("W2", "Feb 10, 2026", "3,610", "0.194%", "0.41%", "0.47×"),
    ("W3", "Mar 10, 2026", "3,687", "0.271%", "0.41%", "0.66×"),
    ("ALL", "Combined",    "10,947","0.237%", "0.41%", "0.58×"),
]
wave_widths = [0.55, 1.35, 1.1, 1.0, 0.8, 0.73]
wy = 3.18
for wi, row in enumerate(wave_data):
    bg = NAVY if wi == 0 else (LGRAY if wi % 2 == 0 else WHITE)
    tc = WHITE if wi == 0 else (GREEN if wi == 4 else DGRAY)
    bld = wi == 0 or wi == 4
    wx = 6.5
    for val, ww in zip(row, wave_widths):
        add_rect(sl, wx, wy, ww, 0.42, fill=bg,
                 line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
        add_text(sl, val, wx+0.04, wy+0.08, ww-0.08, 0.28,
                 size=9, bold=bld, color=tc, align=PP_ALIGN.CENTER)
        wx += ww
    wy += 0.44

add_text(sl, "All 3 waves consistently below BAU. Statistical noise due to sparse "
         "responders (~10 events/window) but directional signal is robust.",
         6.55, 5.05, 6.3, 0.5, size=8.5, color=DGRAY, italic=True)

# Business interpretation
add_rect(sl, 0.3, 6.6, 12.73, 0.6, fill=RGBColor(0xFF, 0xF3, 0xE0),
         line=GOLD, line_w=Pt(1.0))
add_text(sl, "Business profile: The 'passively creditworthy' prospect — clean credit history, comfortable FICO, "
         "not actively shopping for credit. They receive this offer but already have sufficient credit access "
         "and no financial stress triggering card-seeking behaviour.",
         0.45, 6.65, 12.4, 0.52, size=9.5, color=RGBColor(0x8B, 0x5E, 0x00))


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 6 – SHAP Driver Analysis
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "SHAP Driver Analysis — What Predicts Non-Response",
           "XGBoost + TreeSHAP over 26 Equifax bureau features · ranked by mean |SHAP| score")
footer(sl, 6)

drivers = [
    ("EFX_BC_UTIL_1",     0.0578, "Bankcard utilization (% of limit used)",
     "Low utilization = not credit-hungry. Passively using existing cards."),
    ("EFX_AL_BAL_1",      0.0569, "Total open balance across all trades",
     "High overall debt load correlates with either avoidance or already-served need."),
    ("EFX_AL_TRDAGE_1",   0.0357, "Age of oldest open trade (months)",
     "Longer credit history = established borrower with less urgent card need."),
    ("FICO",              0.0351, "Credit score 300–850",
     "Mid-band FICO (664–730) = prime-adjacent but not actively card-shopping."),
    ("EFX_AL_DQOCURNC_3", 0.0308, "DQ events, all trades, 24–36 month window",
     "Absence of recent delinquency signals credit comfort — no pressure to act."),
    ("EFX_BC_BAL_5",      0.0245, "Bankcard balance (36+ month snapshot)",
     "Persistent moderate balance = steady revolving user, not an acquirer."),
    ("EFX_BC_TRDAGE_1",   0.0224, "Age of oldest open bankcard (months)",
     "Established cardholder; existing product portfolio likely meets needs."),
    ("EFX_AL_TRDCNT_1",   0.0213, "Count of all open trade lines",
     "More existing accounts = less need for another credit product."),
    ("EFX_BC_INQCNT_4",   0.0202, "Bankcard inquiries in last 24 months",
     "Inquiry absence confirms no active credit shopping behaviour."),
    ("EFX_BC_TRDCNT_1",   0.0191, "Count of open bankcards",
     "Already holds multiple cards; marginal utility of another is low."),
]

max_shap = 0.0578

# Draw horizontal bar chart
bx, by = 0.3, 1.25
for i, (feat, shap, label, interp) in enumerate(drivers):
    bar_w = 5.5 * (shap / max_shap)
    y = by + i * 0.56
    add_text(sl, f"{i+1}. {label}", bx, y, 5.8, 0.28, size=8.5, color=NAVY, bold=i < 3)
    add_rect(sl, 6.2, y+0.04, bar_w, 0.22, fill=(BLUE if i < 5 else RGBColor(0x90,0xB8,0xE0)))
    add_text(sl, f"{shap:.4f}", 6.2+bar_w+0.05, y+0.04, 0.7, 0.22,
             size=8, color=DGRAY, bold=False)
    add_text(sl, interp, 6.2, y+0.28, 6.8, 0.26, size=7.5, color=DGRAY, italic=True)

# Interpretation box
add_rect(sl, 0.3, 6.85, 12.73, 0.45, fill=RGBColor(0xE8,0xF0,0xFF),
         line=BLUE, line_w=Pt(0.8))
add_text(sl, "Unifying theme: low credit-appetite. Non-responders are passively credit-satisfied — "
         "established trades, adequate limits, modest utilization, no recent shopping. "
         "They receive the offer but have no trigger to act. Segment rules map directly to top SHAP features.",
         0.45, 6.88, 12.4, 0.4, size=9.5, color=BLUE)


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 7 – Business Impact & ROI
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Business Impact — Suppression Economics",
           "Estimated savings from applying the top 5 segments as suppression rules across future mail waves")
footer(sl, 7)

# Summary impact boxes
for i, (lbl, val, sub, vc) in enumerate([
    ("Records Suppressed\n(per wave)", "33,777", "33.8% of 100k file", RED),
    ("Responses Retained", "321 of 410", "78.3% response retention", GREEN),
    ("Responses Forfeited", "89 of 410", "21.7% — acceptable trade-off", GOLD),
    ("Response Rate\n(post-suppression)", "0.48%", "vs 0.41% pre-suppression", GREEN),
]):
    kpi_box(sl, 0.3 + i*3.2, 1.25, 3.0, 1.45, lbl, val, value_color=vc)
    add_text(sl, sub, 0.35 + i*3.2, 2.58, 2.9, 0.28, size=8.5, color=DGRAY,
             align=PP_ALIGN.CENTER)

# Break-even analysis table
add_rect(sl, 0.3, 3.1, 6.1, 3.6, fill=LGRAY, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 0.3, 3.1, 6.1, 0.4, fill=NAVY)
add_text(sl, "BREAK-EVEN MAILING COST ANALYSIS", 0.35, 3.12, 5.8, 0.36, size=11, bold=True, color=WHITE)

be_rows = [
    ("Cost per piece", "$0.50", "$1.00", "$1.50", "$2.00"),
    ("Savings (33,777 pieces)", "$16,889", "$33,777", "$50,666", "$67,554"),
    ("Responses forfeited", "89", "89", "89", "89"),
    ("Break-even value/resp.", "$190", "$380", "$569", "$759"),
    ("Booking rate ~20% → value/booking", "$950", "$1,900", "$2,847", "$3,796"),
    ("NET (saves > lost revenue?", "✓ YES", "✓ YES", "Depends", "Depends"),
]
wy = 3.55
be_widths = [2.2, 0.95, 0.95, 0.95, 0.95]
for ri, row in enumerate(be_rows):
    bg = (NAVY if ri == 0 else (LGRAY if ri % 2 == 1 else WHITE))
    tc = WHITE if ri == 0 else DGRAY
    bld = ri == 0 or ri == 5
    wx = 0.3
    for ci, (val, ww) in enumerate(zip(row, be_widths)):
        fgc = tc
        if ri == 5 and ci > 0:
            fgc = GREEN if '✓' in val else GOLD
        add_rect(sl, wx, wy, ww, 0.38, fill=bg, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
        add_text(sl, val, wx+0.04, wy+0.06, ww-0.08, 0.28, size=8.5, bold=bld,
                 color=fgc, align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)
        wx += ww
    wy += 0.4

add_text(sl, "* Assumes 20% of respondents book a card. Break-even value per response = Savings ÷ Responses forfeited.",
         0.35, 6.6, 5.8, 0.35, size=7.5, color=DGRAY, italic=True)

# Right panel – recommendation
add_rect(sl, 6.7, 3.1, 6.33, 3.6, fill=WHITE, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 6.7, 3.1, 6.33, 0.4, fill=GREEN)
add_text(sl, "RECOMMENDED ACTION", 6.75, 3.12, 6.1, 0.36, size=11, bold=True, color=WHITE)

recs = [
    ("1.", "Implement Segment #1 & #2 immediately",
     "10,947 + 10,881 records. Strongest lift (0.58×), full wave stability, "
     "clean 2–3 condition rules. Operationalise via pre-screen suppression file."),
    ("2.", "Expand to full top 5 after one wave validation",
     "33,777 total. Monitor actual response rates vs prediction. "
     "Refit quarterly as bureau data refreshes."),
    ("3.", "Investigate high-response inverse",
     "V0 cuts show FICO 554–628, high inquiries, and high utilization "
     "as TARGETING signals. Consider positive targeting alongside suppression."),
]
ry = 3.58
for num, title, body in recs:
    add_rect(sl, 6.75, ry, 0.45, 0.82, fill=GREEN)
    add_text(sl, num, 6.75, ry+0.2, 0.45, 0.4, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 7.25, ry+0.02, 5.6, 0.28, size=9.5, bold=True, color=NAVY)
    add_text(sl, body, 7.25, ry+0.3, 5.65, 0.54, size=8.5, color=DGRAY)
    ry += 1.05


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 8 – Operationalisation & Next Steps
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Operationalisation & Next Steps",
           "From segment rules to suppression file — recommended implementation path")
footer(sl, 8)

add_rect(sl, 0, 1.05, 13.33, 0.06, fill=BLUE)

# Suppression profile box
add_rect(sl, 0.3, 1.25, 6.1, 5.7, fill=LGRAY, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 0.3, 1.25, 6.1, 0.42, fill=RED)
add_text(sl, "SUPPRESSION PROFILE — DO-NOT-MAIL INDICATOR", 0.35, 1.27, 5.8, 0.38,
         size=11, bold=True, color=WHITE)

profile_text = """The typical non-responding prospect matches two or more of:

• FICO 664–730  (near-prime; not credit-deprived, not affluent-prime)

• No bankcard DQ occurrences in 24–36 months (EFX_BC_DQOCURNC_3 = 0)

• No all-trade DQ occurrences in 36+ months (EFX_AL_DQOCURNC_5 = 0)

• Zero recent bankcard inquiries in 6 months (EFX_BC_INQCNT_1 = 0)

• Bankcard utilization 16–35% (EFX_BC_UTIL_1 ∈ [16, 35])

Interpretation: Credit-comfortable, passively engaged with existing
products, not actively seeking new credit. Low urgency, low response
probability regardless of offer attractiveness."""

add_text(sl, profile_text, 0.4, 1.75, 5.85, 4.7, size=10, color=DGRAY)

# Steps
add_rect(sl, 6.7, 1.25, 6.33, 5.7, fill=WHITE, line=RGBColor(0xCC,0xCC,0xCC), line_w=Pt(0.5))
add_rect(sl, 6.7, 1.25, 6.33, 0.42, fill=NAVY)
add_text(sl, "IMPLEMENTATION CHECKLIST", 6.75, 1.27, 6.1, 0.38, size=11, bold=True, color=WHITE)

steps = [
    ("Wave N+1", "Pre-screen suppression",
     "Apply top 2 segment rules as hard suppression filters before address-file selection. "
     "Estimated file reduction: 20,664 records (20.7%)."),
    ("Wave N+1", "Track actual vs forecast",
     "Compare realised response rate in suppressed cohort (should be ~0.24%) "
     "to verify model holds on new data."),
    ("Q+1", "Refresh bureau pull",
     "Re-derive EFX attributes from latest bureau snapshot. Re-run V0–V3 pipeline "
     "on refreshed file. Rules may shift as consumer credit conditions change."),
    ("Q+1", "Expand suppression scope",
     "If validation confirms lift, apply full top-5 union. "
     "Target 33% file reduction with <22% response sacrifice."),
    ("Q+2", "Targeting use case",
     "Build positive targeting model from the high-response signals identified "
     "in V0 cuts (high inquiry, sub-prime FICO, high utilization). "
     "Combine suppression + targeting for full file optimisation."),
]

sy = 1.75
for wave, title, body in steps:
    add_rect(sl, 6.75, sy, 1.1, 0.88, fill=BLUE)
    add_text(sl, wave, 6.75, sy+0.22, 1.1, 0.4, size=9, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 7.92, sy+0.02, 4.95, 0.28, size=9.5, bold=True, color=NAVY)
    add_text(sl, body, 7.92, sy+0.3, 5.0, 0.58, size=8.5, color=DGRAY)
    sy += 1.08


# ═══════════════════════════════════════════════════════════════════════════ #
# SLIDE 9 – Appendix: Full Segment Table & SHAP
# ═══════════════════════════════════════════════════════════════════════════ #
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
header_bar(sl, "Appendix — Full Technical Segment Rules",
           "Exact pysubgroup rule syntax from V1 multi-feature beam search · top 12 results shown")
footer(sl, 9)

full_rules = [
    (1, "EFX_AL_DQOCURNC_5==0.0 AND EFX_BC_DQOCURNC_3==0.0 AND FICO: [664.0:694.0[",
     10947, "10.95%", "0.2375%", "0.579"),
    (2, "EFX_AL_DQOCURNC_1==0.0 AND EFX_BC_INQCNT_1==0.0 AND EFX_BC_UTIL_1: [16.0:35.0[",
     10881, "10.88%", "0.2389%", "0.583"),
    (3, "EFX_BC_DQOCURNC_3==0.0 AND EFX_BC_INQCNT_1==0.0 AND FICO: [664.0:694.0[",
     10616, "10.62%", "0.2543%", "0.620"),
    (4, "EFX_BC_DQOCURNC_5==0.0 AND EFX_BC_INQCNT_1==0.0 AND FICO: [694.0:730.0[",
     13060, "13.06%", "0.2603%", "0.635"),
    (5, "EFX_BC_DQOCURNC_3==0.0 AND FICO: [664.0:694.0[",
     12939, "12.94%", "0.2628%", "0.641"),
    (6, "EFX_AL_BAL_1: [67488.0:135936.0[ AND EFX_AL_DQOCURNC_1==0.0",
     13079, "13.08%", "0.2753%", "0.671"),
    (7, "EFX_AL_BAL_1: [67488.0:135936.0[ AND EFX_BC_DQOCURNC_5==0.0",
     15644, "15.64%", "0.2813%", "0.686"),
    (8, "EFX_BC_INQCNT_1==0.0 AND EFX_BC_UTIL_1: [16.0:35.0[",
     15260, "15.26%", "0.2818%", "0.687"),
    (9, "EFX_BC_INQCNT_1==0.0 AND EFX_BC_TRDAGE_2: [48.0:84.0[",
     15334, "15.33%", "0.2869%", "0.700"),
    (10, "EFX_AL_BAL_1: [67488.0:135936.0[",
     18516, "18.52%", "0.2970%", "0.724"),
    (11, "EFX_AL_BAL_1: [67488.0:135936.0[ AND EFX_BC_INQCNT_1==0.0",
     15134, "15.13%", "0.2973%", "0.725"),
    (12, "EFX_BC_INQCNT_1==0.0 AND EFX_BC_TRDAGE_3: [4.0:15.0[",
     15387, "15.39%", "0.2990%", "0.729"),
]

hdr = ["#", "Rule (exact syntax)", "N", "% File", "Resp %", "Lift"]
hw  = [0.3, 7.5, 0.9, 0.75, 0.78, 0.7]
table_row(sl, hdr, 0.3, 1.25, 0.35, hw, bold=True, bg=NAVY, text_color=WHITE)
ry = 1.6
for ri, (num, rule, n, pct, rr, lift) in enumerate(full_rules):
    bg = LGRAY if ri % 2 == 0 else WHITE
    wx = 0.3
    for val, ww in zip([str(num), rule, f"{n:,}", pct, rr, lift], hw):
        add_rect(sl, wx, ry, ww, 0.35, fill=bg, line=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
        al = PP_ALIGN.LEFT if ww == 7.5 else PP_ALIGN.CENTER
        add_text(sl, val, wx+0.04, ry+0.06, ww-0.08, 0.25, size=7.5,
                 color=(RED if val==lift and ri<5 else DGRAY), align=al)
        wx += ww
    ry += 0.37

add_text(sl, "Segments 1–5 recommended for suppression (highlighted). All show response rate consistently "
         "below BAU across 3 mail waves. Segment 6+ show weaker but still-below-BAU rates.",
         0.3, 6.7, 12.73, 0.5, size=8.5, color=DGRAY, italic=True)


# ── Save ────────────────────────────────────────────────────────────────────
import os
os.makedirs("outputs", exist_ok=True)
out = "outputs/Segment_Discovery_Report_2026-06-26.pptx"
prs.save(out)
print(f"Saved -> {out}")
