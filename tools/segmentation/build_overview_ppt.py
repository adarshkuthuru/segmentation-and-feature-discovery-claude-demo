"""
Build the 2-slide executive overview, styled to match the Zenon 2026 template
(docs/templates/Zenon_2026_Template.pptx): 10x5.625in canvas, navy/gold palette,
white cards with a gold left accent bar, Calibri throughout.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

TEMPLATE = "docs/templates/Zenon_2026_Template.pptx"
OUT = "outputs/segmentation/Segment_Discovery_Overview_2026-07-01.pptx"
FONT = "Calibri"

# ── Zenon brand palette (sampled directly from the template) ────────────────
NAVY  = RGBColor(0x1D, 0x2D, 0x44)   # headings, body text, chip fills
GOLD  = RGBColor(0xC8, 0xA1, 0x5A)   # accent bars, numbered chips, dividers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x71, 0x80, 0x96)   # footer / caption text only
BG    = RGBColor(0xF5, 0xF5, 0xF5)   # page background + stat-chip fill


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


prs = Presentation(TEMPLATE)
while len(prs.slides._sldIdLst) > 0:
    delete_slide(prs, 0)
LAYOUT = prs.slide_layouts[0]  # DEFAULT


# ── helpers ───────────────────────────────────────────────────────────────────
def _zero_margins(tf):
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=10, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, anchor=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    _zero_margins(tf)
    tf.word_wrap = wrap
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = FONT
    return tb

def bullets(slide, lines, l, t, w, h, size=7.5, color=NAVY, gap=2.0):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    _zero_margins(tf)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = f"•  {line}"
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tf

def page_header(slide, title, subtitle):
    txt(slide, title, 0.35, 0.08, 9.3, 0.32, size=16, bold=True, color=NAVY)
    box(slide, 0.35, 0.42, 1.0, 0.035, GOLD)
    txt(slide, subtitle, 0.35, 0.48, 9.3, 0.28, size=9, color=GRAY)

def section_label(slide, text, l, t, w):
    txt(slide, text, l, t, w, 0.20, size=9.5, bold=True, color=NAVY)

def footer(slide, page_num):
    txt(slide, "© 2026 Zenon, LLC. All rights reserved.",
        0.35, 5.25, 4.0, 0.25, size=8, color=GRAY)
    txt(slide, str(page_num), 9.3, 5.25, 0.4, 0.25, size=8, color=GRAY,
        align=PP_ALIGN.RIGHT)

def gold_bar_card(slide, l, t, w, h):
    box(slide, l, t, w, h, WHITE)
    box(slide, l, t, 0.05, h, GOLD)

def callout(slide, text, l=0.35, t=0.80, w=9.3, h=0.34, size=9.5):
    """Highlighted one-line takeaway strip. Vertically centered regardless of
    exact font metrics."""
    gold_bar_card(slide, l, t, w, h)
    txt(slide, text, l + 0.22, t, w - 0.4, h, size=size, bold=True,
        color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

def stat_chip(slide, l, t, w, h, label, value):
    box(slide, l, t, w, h, BG)
    box(slide, l, t, w, 0.03, GOLD)
    txt(slide, label, l + 0.02, t + 0.06, w - 0.04, 0.13, size=6, bold=True,
        color=GRAY, align=PP_ALIGN.CENTER)
    txt(slide, value, l + 0.02, t + 0.19, w - 0.04, 0.21, size=9, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — SOLUTION ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(LAYOUT)
bg(s1)
page_header(s1,
    "Solution Architecture: What's Generic vs What's Yours",
    "Generic components work on any dataset — example-specific facts live in the config file or the prompt")
footer(s1, 1)

callout(s1, "Today: analyst supplies both manually (Path A)   →   "
            "Next: Path B derives the config file itself from a one-line prompt (Slide 2)")

CONTENT_TOP = 1.20
CONTENT_BOT = 4.65

# ── Left: Domain-Agnostic components — one box per item ─────────────────────
LX, LW = 0.35, 5.30
section_label(s1, "DOMAIN-AGNOSTIC  —  reusable for any dataset", LX, CONTENT_TOP, LW)

agnostic_items = [
    ("Agent + Skill",
     "Claude Code runs the segment-discovery skill — orchestrates every stage, "
     "narrates results, confirms the spec with the analyst (Path A & Path B)."),
    ("Tool 1: Single-feature decision tree cuts",
     "tree_cuts.py — fits depth-2 trees on every feature independently; surfaces "
     "BAU and best single-feature thresholds. Always the first stage."),
    ("Tool 2: Multi-feature subgroup search",
     "subgroup_search.py — pysubgroup beam search across thousands of rules; "
     "ranks by WRAcc (size × deviation from BAU). The core value-add."),
    ("Tool 3: Stability validation",
     "stability.py — re-checks each top rule within every time/cohort window; "
     "flags unstable segments. No-op when no time column is set."),
    ("Tool 4: SHAP driver analysis",
     "drivers.py — gradient-boosted model + SHAP; confirms rules are "
     "mechanistically grounded, not data artifacts."),
    ("Orchestrator + PPT builder",
     "run_demo.py chains all stages → outputs/segmentation/REPORT.md. build_ppt.py "
     "generates the branded deck. Both are fully data-agnostic."),
    ("JSON run spec contract",
     "spec.template.json defines what to provide (target, features, sentinels…) "
     "but never which values. Swap the spec, the rest is unchanged."),
]

item_top = CONTENT_TOP + 0.24
GAP = 0.05
item_h = (CONTENT_BOT - item_top - GAP * (len(agnostic_items) - 1)) / len(agnostic_items)
for title, body in agnostic_items:
    gold_bar_card(s1, LX, item_top, LW, item_h)
    txt(s1, title, LX + 0.16, item_top + 0.04, LW - 0.32, 0.15, size=7.6, bold=True, color=NAVY)
    txt(s1, body,  LX + 0.16, item_top + 0.20, LW - 0.32, item_h - 0.26, size=6.4, color=NAVY)
    item_top += item_h + GAP

# ── Right: Example-Specific inputs — exactly 2 boxes ────────────────────────
RX, RW = 5.95, 3.70
section_label(s1, "EXAMPLE-SPECIFIC  —  provided by the analyst", RX, CONTENT_TOP, RW)

card_top = CONTENT_TOP + 0.24
card_h = (CONTENT_BOT - card_top - 0.12) / 2

def config_card(top, heading, tag, lines):
    gold_bar_card(s1, RX, top, RW, card_h)
    txt(s1, heading, RX + 0.16, top + 0.06, RW - 0.32, 0.20, size=10, bold=True, color=NAVY)
    txt(s1, tag, RX + 0.16, top + 0.28, RW - 0.32, 0.14, size=6.5, italic=True, color=GRAY)
    bullets(s1, lines, RX + 0.16, top + 0.46, RW - 0.32, card_h - 0.52, size=6.8, gap=1.2)

config_card(card_top, "Config File  —  config.json", "structured & technical", [
    "Data file + dictionary (any CSV/Excel)",
    "Target derivation rule (e.g. responded = 1 where an ID is not null)",
    "Feature set (e.g. 26 EFX attributes) + missing/sentinel codes",
    "Leakage exclusions + time/stability column for validation",
])

card2_top = card_top + card_h + 0.12
config_card(card2_top, "Prompt  —  instructions to Claude", "business & conversational", [
    "Business goal & direction (suppress vs target)",
    "Domain context the analyst already knows",
    "What “good” looks like / review criteria",
    "Ad-hoc guidance not worth encoding in JSON",
])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PATH B: NEXT STEPS (3 consolidated capabilities)
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(LAYOUT)
bg(s2)
page_header(s2,
    "Next Steps: Path B — Cold-Start Without Manual Config",
    "Three capabilities so the agent builds the config file itself, learns over time, and stays within cost bounds")
footer(s2, 2)

callout(s2, "End state: prompt only → agent inspects, researches, drafts config.json, "
            "confirms, and runs")

steps = [
    (1, "Automate the Pipeline Before Path A",
     "multi-agent data prep",
     [
        "Domain Research Agent — web-searches the problem type for standard/missing features",
        "Data Collection Agent — pulls from multiple source files or systems automatically",
        "Data Cleaning Agent — parses the dictionary, detects sentinels/dtypes/leakage",
        "Outputs a draft config.json for one-step analyst approval",
     ],
     "manual column review, dictionary reading, target/leakage selection, and config authoring",
     ("Medium", "3–4 wks", "High")),
    (2, "Persistent Memory Across Runs",
     "warm-start every campaign",
     [
        "Remembers accepted / rejected segment rules across campaigns",
        "Remembers dataset-specific quirks — sentinel codes, leakage columns, feature mappings",
        "New runs on the same or similar data warm-start instead of rediscovering",
     ],
     "re-running full discovery and re-teaching the same dataset quirks every campaign",
     ("Low", "1–2 wks", "Medium")),
    (3, "Self-Correcting Loop With Guardrails",
     "bounded autonomy",
     [
        "Agent re-checks assumptions and refines the spec until a preset objective is met "
        "(e.g. min lift, size, stability)",
        "Hard caps on iterations / token spend stop runaway cost",
        "Escalates to the analyst if the objective isn't met within budget",
     ],
     "single-shot manual review; unattended long-running agents with no cost ceiling",
     ("High", "4–6 wks", "High")),
]

GRID_TOP, GRID_BOT = CONTENT_TOP, CONTENT_BOT
COL_GAP = 0.15
CARD_W = (9.30 - 2 * COL_GAP) / 3
CARD_H = GRID_BOT - GRID_TOP
PAD = 0.14

for idx, (num, title, tag, lines, replaces, (effort, timeline, impact)) in enumerate(steps):
    lx = 0.35 + idx * (CARD_W + COL_GAP)
    ty = GRID_TOP

    box(s2, lx, ty, CARD_W, CARD_H, WHITE)

    # header
    hdr_h = 0.62
    box(s2, lx, ty, CARD_W, hdr_h, NAVY)
    box(s2, lx + 0.1, ty + 0.1, 0.32, 0.32, GOLD)
    txt(s2, str(num), lx + 0.1, ty + 0.1, 0.32, 0.32, size=11, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s2, title, lx + 0.52, ty + 0.08, CARD_W - 0.64, 0.28, size=10.2, bold=True, color=WHITE)
    txt(s2, tag,   lx + 0.52, ty + 0.35, CARD_W - 0.64, 0.2, size=7.5, italic=True, color=GOLD)

    # body bullets
    body_top = ty + hdr_h + 0.08
    bullets_h = 1.50
    bullets(s2, lines, lx + PAD, body_top, CARD_W - 2 * PAD, bullets_h, size=8, gap=3)

    # effort / timeline / impact stat chips
    stat_top = body_top + bullets_h + 0.08
    stat_h = 0.42
    chip_gap = 0.06
    chip_w = (CARD_W - 2 * PAD - 2 * chip_gap) / 3
    for i, (label, value) in enumerate([("EFFORT", effort), ("TIMELINE", timeline), ("IMPACT", impact)]):
        stat_chip(s2, lx + PAD + i * (chip_w + chip_gap), stat_top, chip_w, stat_h, label, value)

    # replaces footer
    div_y = stat_top + stat_h + 0.08
    box(s2, lx + PAD, div_y, CARD_W - 2 * PAD, 0.014, GOLD)
    txt(s2, f"Replaces: {replaces}", lx + PAD, div_y + 0.06, CARD_W - 2 * PAD,
        CARD_H - (div_y - ty) - 0.08, size=6.8, italic=True, color=GRAY)

# ── Bottom summary strip (uses the gap between the grid and the page footer) ─
callout(s2, "Total investment: ~8–12 weeks, phased   →   manual setup drops from "
            "~1 day to minutes per dataset, with cost bounded by design",
        t=4.70, h=0.34, size=9)

prs.save(OUT)
print(f"Saved -> {OUT}")
