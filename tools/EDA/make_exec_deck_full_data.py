"""
Executive deck for the Capital One direct-mail EDA — all 5 phases covered.
Numbers taken from outputs/EDA/Full_Data_eda_report.md (not recomputed).
Charts reused from outputs/EDA/figures/.
Built on docs/templates/Zenon_2026_Template.pptx for visual consistency.
Output: outputs/EDA/Full_Data_executive_summary.pptx
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

BASE = Path(__file__).resolve().parent
ROOT = BASE.parent.parent
OUT_DIR = ROOT / "outputs" / "EDA"
FIG_DIR = OUT_DIR / "figures"
TEMPLATE = ROOT / "docs" / "templates" / "Zenon_2026_Template.pptx"
DECK = OUT_DIR / "Full_Data_executive_summary.pptx"

prs = Presentation(str(TEMPLATE))
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
SLIDE_W_IN = SLIDE_W / 914400
SLIDE_H_IN = SLIDE_H / 914400
MARGIN = 0.35            # horizontal safety margin from the slide edge
BOTTOM_MARGIN = 0.25      # vertical safety margin above the slide's bottom edge
TOP_MARGIN = 0.22         # vertical safety margin below the slide's top edge
GAP_S = 0.08              # small gap: phase tag -> title
GAP_M = 0.15              # standard gap between any two stacked elements
PHASE_H = 0.24
TITLE_H = 0.55
DEFAULT_LAYOUT = prs.slide_layouts[0]

# The template ships with its own pre-built sample slides (title, "About Zenon",
# generic topic) full of placeholder text. Drop them before appending real content
# so the deck doesn't open on unfilled boilerplate.
_xml_slides = prs.slides._sldIdLst
for _sld_id in list(_xml_slides):
    prs.part.drop_rel(_sld_id.rId)
    _xml_slides.remove(_sld_id)

# Brand palette
NAVY  = RGBColor(0x1F, 0x2D, 0x50)
BLUE  = RGBColor(0x4C, 0x72, 0xB0)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
AMBER = RGBColor(0xC8, 0x8A, 0x00)
RED   = RGBColor(0xB0, 0x3A, 0x2E)
GREY  = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xC9, 0xD4, 0xE8)
LTGREY = RGBColor(0xF2, 0xF4, 0xF8)
DKGREY = RGBColor(0x33, 0x33, 0x33)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(DEFAULT_LAYOUT)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def tb(slide, left, top, w, h):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    # Shrink font to fit rather than letting wrapped text silently overflow the box.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    # Internal padding so text never sits flush against the box edge.
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    return tf


def slide_header(s, phase_label, title_text, title_size=22):
    """Places the phase tag, title, and divider with a consistent, non-overlapping
    vertical rhythm, and returns content_top — the y-coordinate every element below
    the header (bullets, images, captions) should start at."""
    phase_tag(s, phase_label, top=TOP_MARGIN)
    title_top = TOP_MARGIN + PHASE_H + GAP_S
    slide_title(s, title_text, top=title_top, size=title_size)
    divider_top = title_top + TITLE_H + GAP_M
    divider_line(s, divider_top)
    return divider_top + GAP_M


def slide_title(slide, text, top=0.28, color=NAVY, size=22):
    tf = tb(slide, 0.45, top, 9.1, 0.6)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color


def phase_tag(slide, label, top=0.28):
    """Small phase label above the title."""
    tf = tb(slide, 0.45, top, 9.1, 0.3)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.size = Pt(9)
    r.font.bold = False
    r.font.color.rgb = BLUE


def body_tf(slide, left=0.55, top=1.0, w=9.0, h=4.2):
    tf = tb(slide, left, top, w, h)
    return tf


def bullet(tf, text, size=14, color=DKGREY, bold=False, indent=0, space_after=6):
    p = tf.add_paragraph()
    r = p.add_run()
    prefix = "    " * indent + ("• " if indent == 0 else "– ")
    r.text = prefix + text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    p.space_after = Pt(space_after)
    return p


def highlight_bullet(tf, text, size=14, color=NAVY, bold=True):
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "→  " + text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    p.space_after = Pt(8)
    return p


def caption(slide, text, top, left=0.45, w=9.1):
    tf = tb(slide, left, top, w, 0.35)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = GREY


def img(slide, fname, left, top, height):
    """Legacy height-only placer, kept for callers that already fit within bounds."""
    path = FIG_DIR / fname
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), height=Inches(height))


def img_fit(slide, fname, left, top, max_w, max_h):
    """Fit an image into an explicit (max_w, max_h) box, scaling by its real aspect
    ratio so it never overflows the slide edge. Returns the actual height used so
    the caller can stack the next element below it."""
    path = FIG_DIR / fname
    if not path.exists():
        return 0
    w_px, h_px = Image.open(path).size
    ar = w_px / h_px
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    left_adj = left + (max_w - w) / 2
    slide.shapes.add_picture(str(path), Inches(left_adj), Inches(top), width=Inches(w), height=Inches(h))
    return h


def divider_line(slide, top, color=LIGHT):
    from pptx.util import Emu
    line = slide.shapes.add_shape(1, Inches(0.45), Inches(top), Inches(9.1), Emu(18000))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def stat_box(slide, left, top, w, h, value, label, val_color=NAVY, bg=LTGREY):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = LIGHT
    tf2 = box.text_frame
    tf2.word_wrap = True
    p1 = tf2.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = val_color
    p2 = tf2.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = GREY


# ===========================================================================
# SLIDE 1 — Title
# ===========================================================================
s = add_slide()
set_bg(s, NAVY)

tf = tb(s, 0.7, 1.1, 8.6, 1.1)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Direct-Mail Prospect Data"
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = WHITE

p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "EDA Readiness Review — All 5 Phases"
r2.font.size = Pt(20)
r2.font.color.rgb = LIGHT

tf2 = tb(s, 0.7, 2.4, 8.6, 0.9)
p3 = tf2.paragraphs[0]
r3 = p3.add_run()
r3.text = "100,000 mailed prospects  •  79 columns  •  3 mailing waves  •  July 2026"
r3.font.size = Pt(13)
r3.font.color.rgb = RGBColor(0x9A, 0xA8, 0xC4)

tf3 = tb(s, 0.7, 3.5, 8.6, 1.6)
for txt, col in [
    ("Phase 1 — Business & Data Understanding", LIGHT),
    ("Phase 2 — Data Quality Checks", LIGHT),
    ("Phase 3 — Univariate & Distribution Analysis", LIGHT),
    ("Phase 4 — Relationship & Target Analysis", LIGHT),
    ("Phase 5 — Time-Based Trends", LIGHT),
]:
    p = tf3.add_paragraph() if txt != "Phase 1 — Business & Data Understanding" else tf3.paragraphs[0]
    r = p.add_run()
    r.text = "  " + txt
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(0x9A, 0xA8, 0xC4)
    p.space_after = Pt(2)


# ===========================================================================
# SLIDE 2 — Phase 1: What this data is
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 1 — BUSINESS & DATA UNDERSTANDING", "One row = one mailed prospect in the campaign")

tf = body_tf(s, top=content_top, h=2.3)
bullet(tf, "100,000 prospects mailed across 3 waves (Jan, Feb, Mar 2026) — the complete campaign file, not a sample.", size=13, bold=True, color=NAVY)
bullet(tf, "79 columns per prospect: credit scores, income, and credit-history snapshots from three bureaus (Equifax, Experian, TransUnion) across multiple product types and time windows.", size=13)
bullet(tf, "Column naming convention: {BUREAU}_{PRODUCT}_{METRIC}_{WINDOW} — e.g. EFX_BC_UTIL_1 = Equifax, Bankcard, Utilisation, Window 1.", size=12, color=GREY)
bullet(tf, 'No pre-built "did they respond?" column exists — we derived it: responded = 1 when BCP_APPLICATION_ID is filled in (an application was submitted).', size=13)
bullet(tf, "Column definitions sourced from the BEST Data Dictionary (847 entries) — no guesswork on field meaning.", size=12, color=GREY)

stat_top = content_top + 2.3 + GAP_M
stat_box(s, 0.45, stat_top, 1.8, 1.3, "100,000", "Prospects mailed")
stat_box(s, 2.35, stat_top, 1.8, 1.3, "79", "Columns")
stat_box(s, 4.25, stat_top, 1.8, 1.3, "3", "Mailing waves")
stat_box(s, 6.15, stat_top, 1.8, 1.3, "410", "Responded (derived target)")
stat_box(s, 8.05, stat_top, 1.8, 1.3, "84", "Booked a card")


# ===========================================================================
# SLIDE 3 — Phase 2 (part 1): Data quality
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 2 — DATA QUALITY CHECKS", "Clean file — hidden sentinel codes needed correcting")

tf = body_tf(s, top=content_top, h=1.1)
bullet(tf, "Zero full-row duplicates — no double-counting of prospects.", size=13, color=GREEN, bold=True)
bullet(tf, "No type mismatches, no broken columns, no out-of-range FICO or utilisation values.", size=13, color=GREEN, bold=True)
bullet(tf, "Date columns (drop date, application date) parse cleanly — no future-dated records.", size=13, color=GREEN, bold=True)

divider2_top = content_top + 1.1 + GAP_M
divider_line(s, divider2_top)

tf2 = body_tf(s, top=divider2_top + GAP_M, h=1.9)
bullet(tf2, 'Credit-bureau fields use sentinel codes (e.g. 999, 9,999,999) to mean "no trade on file" — not actual large values.', size=13, color=AMBER, bold=True)
bullet(tf2, "239,312 sentinel-coded cells across 69 columns were converted to true blanks — without this step they would look like extreme data points.", size=13, color=AMBER)
bullet(tf2, "74 of 79 columns have some missing values after this correction — expected and by design (not every prospect has every trade type).", size=13)
bullet(tf2, "Only 3 columns are >50% blank: BCP_APPLICATION_ID, BCP_ACCOUNT_ID, BCP_APPLICATION_RECEIVED_DATE — all leakage fields excluded from modeling anyway.", size=12, color=GREY)


# ===========================================================================
# SLIDE 4 — Phase 2 (part 2): Bureau redundancy + missingness matrix
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 2 — DATA QUALITY CHECKS (CONTINUED)", "Three bureaus, one signal — largely redundant data")

tf = body_tf(s, left=0.45, top=content_top, w=4.6, h=3.3)
bullet(tf, "23 metric families each appear three times — once from Equifax (EFX), once from Experian (EXP), once from TransUnion (TRU).", size=12)
bullet(tf, "12 of 23 families: near-perfect agreement (avg r ≥ 0.95). Safe to collapse to one representative column.", size=12, color=GREEN, bold=True)
bullet(tf, "11 of 23 families: partial agreement (avg r 0.63–0.93) — especially delinquency-occurrence fields. Bureau disagreement may carry real signal; review before collapsing.", size=12, color=AMBER, bold=True)
bullet(tf, "Missingness is systematic: if a prospect has no trade on file at Equifax, they typically have none at Experian or TransUnion either (dark horizontal bands in the chart).", size=12)

img_left = 5.2
img_h = img_fit(s, "missingness_matrix.png", left=img_left, top=content_top, max_w=SLIDE_W_IN - MARGIN - img_left, max_h=3.4)
caption(s, "Dark bands = missing values across all three bureau copies for the same prospect.", top=content_top + img_h + GAP_M, left=img_left, w=SLIDE_W_IN - MARGIN - img_left)


# ===========================================================================
# SLIDE 5 — Phase 3: Distributions
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 3 — UNIVARIATE & DISTRIBUTION ANALYSIS", "Credit scores are normal; balances and income skew right")

img_h = img_fit(s, "histograms.png", left=0.3, top=content_top, max_w=SLIDE_W_IN - MARGIN - 0.3, max_h=2.0)
cap_top = content_top + img_h + GAP_M
caption(s, "Each bar chart shows how values are spread for one key field (sentinel-cleaned; blanks excluded).", top=cap_top)

tf = body_tf(s, left=0.45, top=cap_top + 0.35 + GAP_M, w=9.1, h=1.0)
bullet(tf, "FICO (credit score): bell-shaped, mean 679, range 434–850. No entry errors.", size=12, color=GREEN)
bullet(tf, "Income & balances: strong right skew (skew 3–13). A small fraction of prospects carry very high balances. Log transform recommended before linear modeling.", size=12, color=AMBER)
bullet(tf, "Delinquency fields: zero-inflated — ~78% of prospects have zero late payments; ~22% of values are flagged as outliers by standard IQR rules (all legitimate, not errors).", size=12, color=AMBER)


# ===========================================================================
# SLIDE 6 — Phase 4a: Target deep-dive
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 4 — RELATIONSHIP & TARGET ANALYSIS", "Response is rare: 243 non-responders per 1 responder")

tf = body_tf(s, left=0.45, top=content_top, w=4.9, h=3.3)
bullet(tf, "0.41% response rate: 410 of 100,000 prospects submitted an application.", size=13, bold=True, color=RED)
bullet(tf, "0.084% booking rate: only 84 of 100,000 opened a card — 4.9× rarer than even responding.", size=13, bold=True, color=RED)
bullet(tf, "Imbalance ratio: 243:1 (non-responders to responders).", size=13)
bullet(tf, 'Standard accuracy is meaningless here — a model that predicts "never responds" for everyone would be 99.59% accurate but completely useless.', size=13, color=AMBER)
bullet(tf, "Correct approach: measure success using lift/ranking (how many times better than random?) and use class-weighted or resampling methods during training.", size=13)

img_left = 5.5
img_h = img_fit(s, "target_balance.png", left=img_left, top=content_top, max_w=SLIDE_W_IN - MARGIN - img_left, max_h=3.4)
caption(s, "Log scale needed to show both bars — the gap between 99,590 and 410 is that large.", top=content_top + img_h + GAP_M, left=img_left, w=SLIDE_W_IN - MARGIN - img_left)


# ===========================================================================
# SLIDE 7 — Phase 4b: Feature vs. target
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 4 — RELATIONSHIP & TARGET ANALYSIS (CONTINUED)", "Responders carry more debt and shop more actively")

img_h = img_fit(s, "feature_by_target.png", left=0.3, top=content_top, max_w=SLIDE_W_IN - MARGIN - 0.3, max_h=2.0)
cap_top = content_top + img_h + GAP_M
caption(s, "Box plots comparing key fields for responders (right) vs. non-responders (left); extreme values trimmed for scale.", top=cap_top)

tf = body_tf(s, left=0.45, top=cap_top + 0.35 + GAP_M, w=9.1, h=1.05)
bullet(tf, "Responders show higher credit utilisation (+4.1 percentage points) and more recent credit inquiries — both signs of active credit-market participation.", size=12, color=GREEN, bold=True)
bullet(tf, "Responders have slightly lower FICO (−3.6 pts) and lower income (−$4,375) than non-responders — the differences are subtle relative to within-group spread.", size=12, color=GREY)
bullet(tf, "Takeaway: utilisation and inquiry activity are more promising early signals than credit score alone.", size=12, color=NAVY, bold=True)


# ===========================================================================
# SLIDE 8 — Phase 4c: Correlation / feature relationships
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 4 — RELATIONSHIP & TARGET ANALYSIS (CONTINUED)", "Each metric family adds unique, non-redundant signal")

tf = body_tf(s, left=0.45, top=content_top, w=4.6, h=3.3)
bullet(tf, "After removing the known bureau triplet repetition, the remaining metric families are essentially unrelated to each other.", size=13)
bullet(tf, "Strongest cross-family correlation: r = 0.011 (near-zero) — balance, age, utilisation, inquiry count, and delinquency each capture a genuinely different dimension of credit behaviour.", size=13, color=GREEN, bold=True)
bullet(tf, "Candidate engineered features:", size=13, bold=True, color=NAVY)
bullet(tf, "Cross-bureau average per family (fills gaps when one bureau is missing)", size=12, indent=1, color=GREY)
bullet(tf, "Has-credit-file flag (thin-file vs. established credit)", size=12, indent=1, color=GREY)
bullet(tf, "Utilisation-to-income ratio (affordability proxy)", size=12, indent=1, color=GREY)
bullet(tf, "Inquiry trend across time windows (credit-shopping acceleration)", size=12, indent=1, color=GREY)

img_left = 5.2
img_h = img_fit(s, "corr_heatmap.png", left=img_left, top=content_top, max_w=SLIDE_W_IN - MARGIN - img_left, max_h=3.4)
caption(s, "One representative per family (bureau siblings excluded). Near-zero correlations confirm independent signal per family.", top=content_top + img_h + GAP_M, left=img_left, w=SLIDE_W_IN - MARGIN - img_left)


# ===========================================================================
# SLIDE 9 — Phase 5: Campaign wave analysis
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "PHASE 5 — TIME-BASED ANALYSIS", "Response rates are stable across all three mailing waves")

img_h = img_fit(s, "waves_trend.png", left=1.0, top=content_top, max_w=SLIDE_W_IN - MARGIN - 1.0, max_h=2.3)
cap_top = content_top + img_h + GAP_M
caption(s, "Bars = volume mailed (left axis). Red line = response rate % (right axis). All three waves near-identical.", top=cap_top)

tf = body_tf(s, left=0.45, top=cap_top + 0.35 + GAP_M, w=9.1, h=0.95)
bullet(tf, "Jan 6: 33,143 mailed → 138 responded (0.416%)    Feb 10: 33,481 mailed → 136 responded (0.406%)    Mar 10: 33,376 mailed → 136 responded (0.407%)", size=12, color=DKGREY)
bullet(tf, "Tight range of 0.406%–0.416% across 3 months — consistent list quality and offer. Time-series deep-dive not applicable (only 3 data points).", size=12, color=GREEN, bold=False)


# ===========================================================================
# SLIDE 10 — Issues to address before modeling
# ===========================================================================
s = add_slide()
content_top = slide_header(s, "SYNTHESIS — PRE-MODELING CHECKLIST", "Seven things to resolve before building a response model")

tf = body_tf(s, top=content_top, h=3.8)
bullet(tf, "Remove 7 leakage columns (BCP_APPLICATION_ID, BCP_ACCOUNT_ID, BCP_APPLICATION_RECEIVED_DATE, SOLICITATION_ID, PV, SRC_ID, ACCT_ID) — these only exist after a prospect has already responded and would give the model the answer.", size=12, color=RED, bold=True)
bullet(tf, 'Decide how to handle "no trade on file" codes — keep as blank, add a has-trade flag, or impute. This analysis blanked them out (correct for stats); production models need a deliberate strategy.', size=12, color=AMBER, bold=True)
bullet(tf, "Collapse the 12 high-agreement bureau triplets (r ≥ 0.95) into one column each — reduces 36 redundant columns to 12.", size=12, color=AMBER)
bullet(tf, "Review and test the 11 lower-agreement bureau triplets (r 0.63–0.93) before collapsing — the delinquency fields in particular may contain bureau-specific signal.", size=12, color=AMBER)
bullet(tf, "Apply log transform to income and balance fields before any linear model — their extreme right skew distorts distances and coefficients.", size=12)
bullet(tf, "Plan the model for extreme rarity (243:1) — use lift/PR-AUC as the success metric, not accuracy; use stratified cross-validation.", size=12, color=RED, bold=True)
bullet(tf, "Engineer cross-bureau averages, a has-credit-file flag, and recency-window trend features to boost signal before model training.", size=12)


# ===========================================================================
# SLIDE 11 — Recommended next steps (closing, dark)
# ===========================================================================
s = add_slide()
set_bg(s, NAVY)

tf = tb(s, 0.7, 0.4, 8.6, 0.75)
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Recommended next steps"
r.font.size = Pt(26)
r.font.bold = True
r.font.color.rgb = WHITE

tf2 = tb(s, 0.9, 1.3, 8.4, 3.8)
steps = [
    ("Clean & simplify first",
     "Apply the sentinel-handling strategy and collapse the 12 safe bureau triplets. This alone reduces the feature space from 71 to ~47 cleaner columns."),
    ("Build a rare-event propensity model",
     "Train a response model using stratified cross-validation and PR-AUC. Evaluate with a lift chart — target the top decile for future mailings."),
    ("Explore the delinquency signal",
     "The 11 lower-agreement bureau triplets (especially DQOCURNC fields) may reveal bureau-specific late-payment patterns. Test whether keeping all three adds lift."),
    ("Check with the campaign team on wave stability",
     "The 0.406%–0.416% wave range is tight but worth a quick conversation — confirm it reflects natural variation, not a list-pull or timing artefact."),
]
for title, detail in steps:
    p = tf2.add_paragraph() if title != steps[0][0] else tf2.paragraphs[0]
    r = p.add_run()
    r.text = f"→  {title}: "
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = LIGHT
    r2 = p.add_run()
    r2.text = detail
    r2.font.size = Pt(13)
    r2.font.bold = False
    r2.font.color.rgb = RGBColor(0xC0, 0xCA, 0xDC)
    p.space_after = Pt(12)

tf3 = tb(s, 0.7, 5.1, 8.6, 0.4)
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Data is structurally clean and ready — complete the pre-modeling checklist, then proceed to modeling."
r.font.size = Pt(11)
r.font.italic = True
r.font.color.rgb = RGBColor(0x9A, 0xA8, 0xC4)

# ---------------------------------------------------------------------------
prs.save(str(DECK))
print(f"Executive deck written to: {DECK}")
print(f"Slides: {len(prs.slides)}")
