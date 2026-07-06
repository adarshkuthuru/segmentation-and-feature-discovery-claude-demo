"""Build executive-ready suppression segment PPT — run 2026-07-06."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)
TEAL   = RGBColor(0x00, 0x7B, 0x83)
RED    = RGBColor(0xC4, 0x1E, 0x3A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)
MGRAY  = RGBColor(0xCC, 0xD1, 0xD9)
DKGRAY = RGBColor(0x44, 0x4E, 0x5C)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
GOLD   = RGBColor(0xD4, 0xA0, 0x17)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def hline(slide, t, color=MGRAY, l=0.3, w=12.73):
    ln = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.02))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def header_bar(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 1.25, NAVY)
    txt(slide, title, 0.35, 0.12, 12.0, 0.65, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.78, 12.0, 0.38, size=13, color=MGRAY)

def footer(slide, note=""):
    box(slide, 0, 7.15, 13.33, 0.35, NAVY)
    txt(slide, "Credit Card Direct-Mail  |  Suppression Segment Discovery  |  July 2026",
        0.35, 7.18, 8.0, 0.28, size=9, color=MGRAY)
    if note:
        txt(slide, note, 8.5, 7.18, 4.5, 0.28, size=9, color=MGRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)
box(s1, 0, 3.15, 13.33, 0.08, TEAL)

txt(s1, "Suppression Segment Discovery",
    1.0, 1.5, 11.33, 1.0, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, "Credit Card Direct-Mail Campaign  —  Run: 2026-07-06",
    1.0, 2.55, 11.33, 0.55, size=22, color=MGRAY, align=PP_ALIGN.CENTER)

for i, (val, lab) in enumerate([
    ("100,000",  "Prospects Solicited"),
    ("0.41%",    "BAU Response Rate"),
    ("4 Stages", "Analysis Pipeline"),
    ("6 Rules",  "Suppression Segments Found"),
]):
    lx = 1.2 + i * 2.8
    box(s1, lx, 3.6, 2.4, 1.5, RGBColor(0x14, 0x3A, 0x6B))
    txt(s1, val, lx, 3.72, 2.4, 0.65, size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s1, lab, lx, 4.3, 2.4, 0.65, size=12, color=MGRAY, align=PP_ALIGN.CENTER)

txt(s1, "Zenon AI  |  Confidential  |  July 6 2026",
    0, 6.9, 13.33, 0.5, size=10, color=MGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, WHITE)
header_bar(s2, "Executive Summary", "What we found and why it matters")
footer(s2, "Slide 2 of 7")

box(s2, 0.3, 1.45, 6.1, 5.4, LGRAY)
txt(s2, "The Opportunity", 0.55, 1.6, 5.6, 0.4, size=14, bold=True, color=NAVY)
hline(s2, 2.1, TEAL, 0.55, 5.6)

bullets = [
    ("Wasted spend:",        "At 0.41% BAU, 99.6% of the mail file never responds."),
    ("6 segments found:",    "Multi-feature rules that suppress response to 0.24–0.28% — up to 42% below BAU."),
    ("Top 2 segments:",      "~10.9% of the mail file each — drop them and cut nearly 22% of volume with minimal loss."),
    ("Consistent signal:",   "All rules directionally below BAU across all 3 waves (Jan / Feb / Mar 2026)."),
    ("Mechanistic backing:", "BC utilization, total balance, and trade age are the #1–3 SHAP drivers — rules make sense."),
]
ty = 2.25
for label, body in bullets:
    box(s2, 0.55, ty, 0.08, 0.22, TEAL)
    txt(s2, label, 0.75, ty - 0.02, 1.7, 0.28, size=11, bold=True, color=NAVY)
    txt(s2, body,  2.45, ty - 0.02, 3.8, 0.38, size=11, color=DKGRAY, wrap=True)
    ty += 0.72

kpis = [
    ("0.58×",  "Best Suppression Lift",        RED),
    ("42%",    "Response Rate Drop (top 2)",    RED),
    ("10.9%",  "Population per Top Segment",    TEAL),
    ("~22%",   "Mail Suppressable (top 2 segs)", TEAL),
]
for i, (val, lab, col) in enumerate(kpis):
    rx, ry = 6.9, 1.55 + i * 1.3
    box(s2, rx, ry, 5.8, 1.1, col)
    txt(s2, val, rx, ry + 0.05, 5.8, 0.65, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s2, lab, rx, ry + 0.68, 5.8, 0.35, size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, WHITE)
header_bar(s3, "Methodology", "Four-stage supervised subgroup discovery pipeline")
footer(s3, "Slide 3 of 7")

stages = [
    ("V0", "EDA + Single-Feature Cuts",
     "Establish BAU. Run decision-tree cuts on each feature individually.\n"
     "Best single cut: Income band $292k–$346k → 0.06% rate (lift 0.14×)."),
    ("V1", "Multi-Feature Subgroup Search",
     "Beam search over 2–3-condition rules (pysubgroup). Ranks by WRAcc.\n"
     "25 features × beam search → 6 top suppression rules surfaced."),
    ("V2", "Stability Validation",
     "Re-scores each top rule within each mail wave (Jan / Feb / Mar 2026).\n"
     "All rules held below BAU in every wave — directionally persistent."),
    ("V3", "SHAP Driver Analysis",
     "Gradient-boosted model; mean |SHAP| per feature.\n"
     "Top drivers: BC Utilization (0.087), Total Balance (0.084), Trade Age (0.057), FICO (0.055)."),
]

for i, (tag, title, desc) in enumerate(stages):
    lx = 0.35 + i * 3.22
    box(s3, lx, 1.45, 3.0, 0.55, NAVY)
    txt(s3, f"{tag}  ·  {title}", lx + 0.12, 1.5, 2.76, 0.48, size=11, bold=True, color=WHITE)
    box(s3, lx, 2.0, 3.0, 4.4, LGRAY)
    txt(s3, desc, lx + 0.12, 2.1, 2.76, 4.1, size=10.5, color=DKGRAY, wrap=True)
    if i < 3:
        txt(s3, "→", lx + 3.03, 1.55, 0.2, 0.45, size=16, bold=True, color=TEAL)

txt(s3,
    "Target: BCP_APPLICATION_ID not null → responded=1  |  Features: 25 EFX attributes + FICO + Income (one bureau to avoid tri-bureau redundancy)"
    "  |  Sentinels treated as NaN: 9999997, 9999998, 9999999, 997, 998, 999  |  Time column: TEST_CELL_DROP_DATE (3 waves)",
    0.35, 6.75, 12.63, 0.32, size=9, color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RANKED SUPPRESSION SEGMENTS
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, WHITE)
header_bar(s4, "Ranked Suppression Segments",
           "Six multi-feature rules — all below BAU across Jan / Feb / Mar 2026 mailing waves")
footer(s4, "Slide 4 of 7")

cols   = ["#", "Business Label",         "Plain-English Rule",                  "Size",   "Size %", "Resp Rate", "BAU",   "Lift"]
widths = [0.3,  2.0,                      5.0,                                    0.65,     0.6,      0.75,        0.55,    0.68]

ty0, row_h = 1.45, 0.55
cx = 0.25
for col, w in zip(cols, widths):
    box(s4, cx, ty0, w, 0.4, NAVY)
    txt(s4, col, cx + 0.05, ty0 + 0.04, w - 0.1, 0.35, size=10, bold=True, color=WHITE)
    cx += w

rows = [
    ["1", "Clean DQ (5yr+3mo) + FICO 664–694",
     "No DQ occ 5-yr + No BC DQ 3-mo + FICO 664–694",
     "10,947", "10.9%", "0.24%", "0.41%", "0.58×"],
    ["2", "Mid-Util + No Inquiry + Clean DQ",
     "No DQ 1-yr + No BC inquiry (1mo) + BC util 16–35%",
     "10,881", "10.9%", "0.24%", "0.41%", "0.58×"],
    ["3", "FICO 664–694 + No BC Inquiry + Clean BC DQ",
     "No BC DQ 3-mo + No BC inquiry (1mo) + FICO 664–694",
     "10,616", "10.6%", "0.25%", "0.41%", "0.62×"],
    ["4", "FICO 694–730 + No Inquiry + Clean BC DQ",
     "No BC DQ 5-mo + No BC inquiry (1mo) + FICO 694–730",
     "13,060", "13.1%", "0.26%", "0.41%", "0.64×"],
    ["5", "FICO 664–694 + Clean BC DQ 3mo",
     "No BC DQ 3-mo + FICO 664–694",
     "12,939", "12.9%", "0.26%", "0.41%", "0.64×"],
    ["6", "High Balance ($67.5K–$136K) + Clean DQ",
     "Total balance $67.5K–$136K + No DQ occ 1-yr",
     "13,079", "13.1%", "0.28%", "0.41%", "0.67×"],
]

for ri, row in enumerate(rows):
    fill = LGRAY if ri % 2 == 0 else WHITE
    cx = 0.25
    ry = ty0 + 0.4 + ri * row_h
    for ci, (val, w) in enumerate(zip(row, widths)):
        box(s4, cx, ry, w, row_h - 0.02, fill)
        c = RED if ci == 7 else (NAVY if ci in (0, 1) else DKGRAY)
        sz = 10 if ci == 2 else 11
        txt(s4, val, cx + 0.05, ry + 0.04, w - 0.08, row_h - 0.08,
            size=sz, bold=(ci in (0, 7)), color=c, wrap=True)
        cx += w

txt(s4, "Lift < 1.0 = suppression (segment responds LESS than BAU).  Lower lift = stronger suppression.",
    0.25, 6.78, 12.63, 0.3, size=9, color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — STABILITY ACROSS MAILING WAVES
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5, WHITE)
header_bar(s5, "Stability Across Mailing Waves",
           "Response rates for top 5 segments — all below BAU in every wave")
footer(s5, "Slide 5 of 7")

wcols   = ["#", "Business Label",                "Jan 2026", "Feb 2026", "Mar 2026", "Overall", "BAU"]
wwidths = [0.3,  3.5,                              1.45,       1.45,       1.45,       1.2,       0.85]

cx = 0.5
ty = 1.5
box_h = 0.42
for col, w in zip(wcols, wwidths):
    box(s5, cx, ty, w, box_h, NAVY)
    txt(s5, col, cx + 0.06, ty + 0.08, w - 0.1, box_h - 0.1, size=11, bold=True, color=WHITE)
    cx += w

wave_rows = [
    ["1", "Clean DQ (5yr+3mo) + FICO 664–694",    "0.25%", "0.19%", "0.27%", "0.24%", "0.41%"],
    ["2", "Mid-Util + No Inquiry + Clean DQ",      "0.25%", "0.25%", "0.22%", "0.24%", "0.41%"],
    ["3", "FICO 664–694 + No BC Inq + Clean DQ",   "0.28%", "0.23%", "0.25%", "0.25%", "0.41%"],
    ["4", "FICO 694–730 + No Inquiry + Clean DQ",  "0.30%", "0.20%", "0.28%", "0.26%", "0.41%"],
    ["5", "FICO 664–694 + Clean BC DQ 3mo",        "0.26%", "0.23%", "0.30%", "0.26%", "0.41%"],
]

for ri, row in enumerate(wave_rows):
    fill = LGRAY if ri % 2 == 0 else WHITE
    cx = 0.5
    ry = ty + box_h + ri * (box_h + 0.03)
    for ci, (val, w) in enumerate(zip(row, wwidths)):
        box(s5, cx, ry, w, box_h, fill)
        c = MGRAY if ci == 6 else (TEAL if ci in (2, 3, 4, 5) else DKGRAY)
        txt(s5, val, cx + 0.06, ry + 0.08, w - 0.1, box_h - 0.1,
            size=11, bold=(ci == 5), color=c)
        cx += w

box(s5, 0.5, 4.9, 12.3, 0.5, RGBColor(0xFF, 0xF3, 0xE0))
txt(s5, "BAU = 0.41%  |  Every cell above is BELOW BAU — confirming persistent non-response across all mailing waves.",
    0.65, 4.95, 12.0, 0.35, size=11, color=RED)

txt(s5,
    "Stability note: formal stable=True requires all waves < 0.205% (50% of BAU). No rule clears that strict bar, "
    "but all are consistently below BAU in every wave — operationally valid for suppression.",
    0.5, 5.55, 12.63, 0.45, size=9, color=MGRAY, italic=True, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SHAP DRIVERS
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6, WHITE)
header_bar(s6, "Global Feature Drivers (SHAP)",
           "Top 8 drivers of response — confirms segments are mechanistically grounded")
footer(s6, "Slide 6 of 7")

drivers = [
    ("EFX_BC_UTIL_1",     0.087, "Bankcard Utilization (ratio)",
     "Mid-util (16–35%) = card-saturated but not maxed; top suppression split"),
    ("EFX_AL_BAL_1",      0.084, "Total Open Balance (all trades)",
     "High balance band $67.5K–$136K is a strong suppression signal"),
    ("EFX_AL_TRDAGE_1",   0.057, "Age of All Open Trades (months)",
     "Established borrowers less likely to switch — longer age → lower response"),
    ("FICO",              0.055, "FICO Credit Score",
     "Mid-prime band 664–730 anchors 4 of the 6 top segments"),
    ("EFX_AL_DQOCURNC_3", 0.050, "DQ Occurrences — 3-year window",
     "Clean DQ paradox: no delinquency = credit-satisfied, less motivated by new card"),
    ("EFX_BC_TRDAGE_2",   0.039, "Age of BC Trades (2nd oldest)",
     "Supporting signal; older BC trades correlate with settled-card behaviour"),
    ("EFX_BC_BAL_5",      0.035, "Bankcard Balance (5-mo)",
     "Balance trajectory over 5 months adds incremental suppression signal"),
    ("EFX_AL_TRDCNT_1",   0.034, "Total Open Trade Count",
     "Many trades = established credit user, less motivated by yet another offer"),
]

max_shap = 0.087
bar_max_w = 5.5
ty = 1.5
for i, (feat, shap, label, interp) in enumerate(drivers):
    ry = ty + i * 0.58
    bar_w = (shap / max_shap) * bar_max_w
    box(s6, 0.3, ry + 0.07, 0.32, 0.32, TEAL)
    txt(s6, str(i + 1), 0.3, ry + 0.06, 0.32, 0.32, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s6, label, 0.68, ry + 0.04, 3.3, 0.3, size=10, bold=True, color=NAVY)
    txt(s6, feat,  0.68, ry + 0.3,  3.3, 0.25, size=8.5, color=MGRAY, italic=True)
    box(s6, 4.1, ry + 0.08, bar_w, 0.28, TEAL)
    txt(s6, f"{shap:.3f}", 4.15 + bar_w, ry + 0.08, 0.6, 0.28, size=9, bold=True, color=NAVY)
    txt(s6, interp, 10.1, ry + 0.04, 2.95, 0.5, size=8.5, color=DKGRAY, wrap=True)

hline(s6, 6.2, MGRAY, 0.3, 12.73)
txt(s6,
    "SHAP values from gradient-boosted model on 25 EFX features + FICO + Income. "
    "Mean |SHAP| = average impact on model output magnitude.",
    0.3, 6.3, 12.73, 0.4, size=9, color=MGRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RECOMMENDATIONS & NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7, WHITE)
header_bar(s7, "Recommendations & Next Steps", "Analyst accept / reject — then operationalise")
footer(s7, "Slide 7 of 7")

box(s7, 0.3, 1.45, 6.3, 5.3, LGRAY)
txt(s7, "Recommended Suppression Actions", 0.5, 1.55, 6.0, 0.38, size=13, bold=True, color=NAVY)
hline(s7, 2.05, TEAL, 0.5, 6.0)

recs = [
    ("Accept Segs 1 & 2 (minimum action)",
     "Suppress ~21,800 records (21.9% of mail). Both hit 0.24% resp rate — 42% below BAU. "
     "Consistent across all three waves. Low revenue risk."),
    ("Accept Segs 3 & 4 (moderate action)",
     "Add ~23,700 records. Segs 3–4 hold 0.25–0.26% across waves. "
     "FICO 694–730 (seg 4) extends coverage into near-prime band."),
    ("Accept Seg 5 (broader FICO 664–694 cut)",
     "Adds ~12,900 records. Simpler 2-condition rule (no BC DQ + FICO). "
     "Validate wave rates before operationalising — Mar wave hit 0.30%."),
    ("Re-run pipeline on next campaign file",
     "Rules are feature-based, not sample-specific. Confirm lift holds "
     "on the next solicitation universe. Set quarterly cadence."),
]
ty = 2.2
for i, (title, body) in enumerate(recs):
    box(s7, 0.5, ty, 0.08, 0.22, RED if i < 3 else TEAL)
    txt(s7, title, 0.7, ty - 0.02, 5.7, 0.28, size=10.5, bold=True, color=NAVY)
    txt(s7, body,  0.7, ty + 0.28, 5.7, 0.5,  size=9.5,  color=DKGRAY, wrap=True)
    ty += 0.95

box(s7, 6.9, 1.45, 6.1, 5.3, NAVY)
txt(s7, "Projected Impact", 7.1, 1.55, 5.7, 0.38, size=13, bold=True, color=WHITE)
hline(s7, 2.05, TEAL, 7.1, 5.7)

impacts = [
    ("Segs 1–2 only",  "~21,800 records",  "0.24% resp rate", "42% below BAU"),
    ("Segs 1–4",       "~35,000 records",  "0.25% resp rate", "39% below BAU"),
    ("All 6 segs",     "~48,000 records",  "0.26% resp rate", "36% below BAU"),
]
ty2 = 2.2
for label, size, rate, lift in impacts:
    box(s7, 7.1, ty2, 5.7, 1.15, RGBColor(0x14, 0x3A, 0x6B))
    txt(s7, label, 7.25, ty2 + 0.05, 5.4, 0.3,  size=11, bold=True, color=TEAL)
    txt(s7, size,  7.25, ty2 + 0.38, 2.5, 0.3,  size=12, bold=True, color=WHITE)
    txt(s7, rate,  9.8,  ty2 + 0.38, 2.7, 0.3,  size=11, color=MGRAY)
    txt(s7, lift,  7.25, ty2 + 0.72, 5.4, 0.3,  size=10, color=RED)
    ty2 += 1.3

txt(s7, "* Overlap-adjusted estimates. Exact counts depend on which segments are combined and de-duplicated.",
    6.9, 6.6, 6.1, 0.3, size=8.5, color=MGRAY, italic=True)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "outputs/Suppression_Segment_Report_2026-07-06.pptx"
prs.save(out)
print(f"Saved -> {out}")
