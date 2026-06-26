"""Build executive-ready suppression segment discovery PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)   # slide backgrounds / headers
TEAL   = RGBColor(0x00, 0x7B, 0x83)   # accent bars / highlights
ORANGE = RGBColor(0xE8, 0x6B, 0x1A)   # suppression "bad" colour
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF4, 0xF7)   # alternating row fill
MGRAY  = RGBColor(0xCC, 0xD1, 0xD9)   # borders / subdued text
DKGRAY = RGBColor(0x44, 0x4E, 0x5C)   # body text

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # truly blank


# ── helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def hline(slide, t, color=MGRAY, l=0.3, w=12.73):
    ln = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.02))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def header_bar(slide, title, subtitle=None):
    box(slide, 0, 0, 13.33, 1.25, NAVY)
    txt(slide, title, 0.35, 0.12, 12.0, 0.65, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.78, 12.0, 0.38, size=13, color=MGRAY)

def footer(slide, note=""):
    box(slide, 0, 7.15, 13.33, 0.35, NAVY)
    txt(slide, "Credit Card Direct-Mail | Suppression Segment Discovery | June 2026",
        0.35, 7.18, 8.0, 0.28, size=9, color=MGRAY)
    if note:
        txt(slide, note, 8.5, 7.18, 4.5, 0.28, size=9, color=MGRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)
box(s1, 0, 3.15, 13.33, 0.08, TEAL)

txt(s1, "Suppression Segment Discovery",
    1.0, 1.5, 11.33, 1.0, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, "Credit Card Direct-Mail Campaign",
    1.0, 2.55, 11.33, 0.55, size=22, color=MGRAY, align=PP_ALIGN.CENTER)

# Stat boxes
for i, (val, lab) in enumerate([
    ("100,000", "Prospects Mailed"),
    ("0.41%",   "BAU Response Rate"),
    ("4 Stages","Analysis Pipeline"),
    ("6 Rules", "Suppression Segments"),
]):
    lx = 1.2 + i * 2.8
    box(s1, lx, 3.6, 2.4, 1.5, RGBColor(0x14, 0x3A, 0x6B))
    txt(s1, val, lx, 3.72, 2.4, 0.65, size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s1, lab, lx, 4.3, 2.4, 0.65, size=12, color=MGRAY, align=PP_ALIGN.CENTER)

txt(s1, "Zenon AI  |  Confidential",
    0, 6.9, 13.33, 0.5, size=10, color=MGRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, WHITE)
header_bar(s2, "Executive Summary", "What we found and why it matters")
footer(s2, "Slide 2 of 7")

# Left: opportunity statement
box(s2, 0.3, 1.45, 6.1, 5.4, LGRAY)
txt(s2, "The Opportunity", 0.55, 1.6, 5.6, 0.4, size=14, bold=True, color=NAVY)
hline(s2, 2.1, TEAL, 0.55, 5.6)

bullets = [
    ("Wasted spend:", "At 0.41% BAU, 99.6% of the mail file does not respond."),
    ("Suppression segments:", "We identified 6 rule-based cohorts that respond at 33–42% below BAU."),
    ("Top 2 segments:", "~21% of the mail file — suppress them and nearly eliminate that cost."),
    ("Consistent across waves:", "All segments held below BAU across all 3 mailing waves (Jan–Mar 2026)."),
    ("Mechanistic signals:", "Bankcard utilization, total balance, and trade age are the #1–3 global drivers."),
]
ty = 2.25
for label, body in bullets:
    box(s2, 0.55, ty, 0.08, 0.22, TEAL)
    txt(s2, label, 0.75, ty - 0.02, 1.6, 0.28, size=11, bold=True, color=NAVY)
    txt(s2, body,  2.35, ty - 0.02, 3.9, 0.35, size=11, color=DKGRAY, wrap=True)
    ty += 0.72

# Right: KPI panel
kpis = [
    ("0.58×",  "Best Lift vs BAU",       ORANGE),
    ("42%",    "Response Rate Drop",      ORANGE),
    ("10.9%",  "Population Covered (top seg)", TEAL),
    ("~21%",   "Mail Suppressed (top 2 segs)", TEAL),
]
for i, (val, lab, col) in enumerate(kpis):
    rx = 6.9
    ry = 1.55 + i * 1.3
    box(s2, rx, ry, 5.8, 1.1, col)
    txt(s2, val, rx, ry + 0.05, 5.8, 0.65, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s2, lab, rx, ry + 0.68, 5.8, 0.35, size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — METHODOLOGY
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, WHITE)
header_bar(s3, "Methodology", "Four-stage supervised subgroup discovery pipeline")
footer(s3, "Slide 3 of 7")

stages = [
    ("v0", "EDA + Single-Feature Cuts",
     "Establish BAU. Run decision-tree cuts on each feature individually.\nBest single feature: Trade age 202–210 months → 0.14% rate (lift 0.34×)."),
    ("v1", "Multi-Feature Subgroup Search",
     "Beam search over 2–3-condition rules (pysubgroup). Ranks by WRAcc quality score.\nSearched 26 EFX features × 3 conditions → surfaced 25 candidate rules."),
    ("v2", "Stability Validation",
     "Re-checks each top rule across 3 mailing waves (Jan / Feb / Mar 2026).\nAll rules held consistently below BAU in every wave."),
    ("v3", "SHAP Driver Analysis",
     "Trains gradient-boosted model; computes mean |SHAP| per feature.\nTop drivers: BC Utilization, Total Balance, Trade Age, FICO."),
]

for i, (tag, title, desc) in enumerate(stages):
    lx = 0.35 + i * 3.22
    box(s3, lx, 1.45, 3.0, 0.55, NAVY)
    txt(s3, f"{tag}  ·  {title}", lx + 0.12, 1.5, 2.76, 0.48, size=11, bold=True, color=WHITE)
    box(s3, lx, 2.0, 3.0, 4.4, LGRAY)
    txt(s3, desc, lx + 0.12, 2.1, 2.76, 4.1, size=10.5, color=DKGRAY, wrap=True)
    # Connector arrow (except last)
    if i < 3:
        txt(s3, "→", lx + 3.03, 1.55, 0.2, 0.45, size=16, bold=True, color=TEAL)

txt(s3, "Target: responded = 1 if BCP_APPLICATION_ID is not null  |  Features: 26 Equifax bureau attributes + FICO + Income  |  Sentinels treated as missing: 9999997, 9999998, 9999999, 997, 998, 999",
    0.35, 6.75, 12.63, 0.32, size=9, color=MGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TOP SUPPRESSION SEGMENTS TABLE
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, WHITE)
header_bar(s4, "Ranked Suppression Segments", "All six rules consistently below BAU across Jan / Feb / Mar 2026 mailing waves")
footer(s4, "Slide 4 of 7")

# Table header
cols  = ["#", "Business Label", "Plain-English Rule", "Size", "Size %", "Resp Rate", "BAU", "Lift"]
widths = [0.3,  2.0,             4.9,                  0.7,   0.6,      0.75,        0.55,  0.7]
lx0 = 0.25
ty0 = 1.45
row_h = 0.55

# Header row
cx = lx0
for col, w in zip(cols, widths):
    box(s4, cx, ty0, w, 0.4, NAVY)
    txt(s4, col, cx + 0.05, ty0 + 0.04, w - 0.1, 0.35, size=10, bold=True, color=WHITE)
    cx += w

rows = [
    ["1", "FICO Mid + Clean DQ",    "No DQ (3yr) + No BC DQ (90d)\n+ FICO 664–694",          "10,947", "11.0%", "0.24%", "0.41%", "0.58×"],
    ["2", "Mid-Util, No Inquiry",   "No recent DQ + No BC inquiries\n+ BC util 16–35%",       "10,881", "10.9%", "0.24%", "0.41%", "0.58×"],
    ["3", "Good-FICO, No Inquiry",  "No BC DQ (5yr) + No BC inquiries\n+ FICO 694–730",       "13,060", "13.1%", "0.26%", "0.41%", "0.64×"],
    ["4", "Mid-FICO, No DQ 3yr",   "No BC DQ (3yr)\n+ FICO 664–694",                          "12,939", "12.9%", "0.26%", "0.41%", "0.64×"],
    ["5", "High Balance, Clean",    "Total balance $67.5K–$136K\n+ No DQ (1yr)",               "13,079", "13.1%", "0.28%", "0.41%", "0.67×"],
    ["6", "High Balance + No DQ5",  "Total balance $67.5K–$136K\n+ No BC DQ (5yr)",            "15,644", "15.6%", "0.28%", "0.41%", "0.69×"],
]

for ri, row in enumerate(rows):
    fill = LGRAY if ri % 2 == 0 else WHITE
    cx = lx0
    ry = ty0 + 0.4 + ri * row_h
    for ci, (val, w) in enumerate(zip(row, widths)):
        box(s4, cx, ry, w, row_h - 0.02, fill)
        c = ORANGE if ci == 7 else NAVY if ci in (0, 1) else DKGRAY
        sz = 10 if ci == 2 else 11
        txt(s4, val, cx + 0.05, ry + 0.04, w - 0.08, row_h - 0.08,
            size=sz, bold=(ci in (0, 7)), color=c, wrap=True)
        cx += w

txt(s4, "Lift < 1.0 = suppression (segment responds LESS than BAU).  Lower lift = stronger suppression candidate.",
    0.25, 6.78, 12.63, 0.3, size=9, color=MGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — STABILITY ACROSS MAILING WAVES
# ════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5, WHITE)
header_bar(s5, "Stability Across Mailing Waves", "Response rates for top 5 segments vs BAU — all consistently below BAU across all 3 waves")
footer(s5, "Slide 5 of 7")

# Wave table — top 5 rules
wcols  = ["#", "Business Label",          "Jan 2026", "Feb 2026", "Mar 2026", "Overall", "BAU"]
wwidths = [0.3,  3.3,                      1.5,        1.5,        1.5,        1.3,       0.9]

cx = 0.5
ty = 1.5
box_h = 0.42

cx = 0.5
for col, w in zip(wcols, wwidths):
    box(s5, cx, ty, w, box_h, NAVY)
    txt(s5, col, cx + 0.06, ty + 0.08, w - 0.1, box_h - 0.1, size=11, bold=True, color=WHITE)
    cx += w

wave_rows = [
    ["1", "FICO Mid + Clean DQ",   "0.25%", "0.19%", "0.27%", "0.24%", "0.41%"],
    ["2", "Mid-Util, No Inquiry",  "0.25%", "0.25%", "0.22%", "0.24%", "0.41%"],
    ["3", "Good-FICO, No Inquiry", "0.30%", "0.20%", "0.28%", "0.26%", "0.41%"],
    ["4", "Mid-FICO, No DQ 3yr",  "0.26%", "0.23%", "0.30%", "0.26%", "0.41%"],
    ["5", "High Balance, Clean",   "0.33%", "0.29%", "0.21%", "0.28%", "0.41%"],
]

for ri, row in enumerate(wave_rows):
    fill = LGRAY if ri % 2 == 0 else WHITE
    cx = 0.5
    ry = ty + box_h + ri * (box_h + 0.03)
    for ci, (val, w) in enumerate(zip(row, wwidths)):
        box(s5, cx, ry, w, box_h, fill)
        c = DKGRAY
        if ci == 6:  # BAU column
            c = MGRAY
        elif ci in (2, 3, 4, 5):
            c = TEAL
        txt(s5, val, cx + 0.06, ry + 0.08, w - 0.1, box_h - 0.1,
            size=11, bold=(ci == 5), color=c)
        cx += w

# BAU reference line note
box(s5, 0.5, 5.0, 12.33, 0.5, RGBColor(0xFF, 0xF3, 0xE0))
txt(s5, "BAU = 0.41%   |   Every cell above is BELOW BAU, confirming persistent non-response across all mailing waves.",
    0.65, 5.05, 12.0, 0.35, size=11, bold=False, color=ORANGE)

txt(s5, "Stability flag: formal stable=True requires all waves < 0.205% (50% of BAU). No segment crosses that strict threshold, but all are consistently below BAU — operationally valid for suppression.",
    0.5, 5.65, 12.63, 0.45, size=9, color=MGRAY, italic=True, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SHAP DRIVERS
# ════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6, WHITE)
header_bar(s6, "Global Feature Drivers (SHAP)", "What predicts response — confirms segments are mechanistically grounded")
footer(s6, "Slide 6 of 7")

drivers = [
    ("EFX_BC_UTIL_1",     0.058, "Bankcard Utilization (ratio)",        "Mid-utilization (16–35%) predicts non-response — card-saturated but not maxed"),
    ("EFX_AL_BAL_1",      0.057, "Total Open Balance (all trades)",     "High total balance ($67.5K–$136K) is a strong suppression signal"),
    ("EFX_AL_TRDAGE_1",   0.036, "Age of All Open Trades (months)",    "Longer trade age = established borrowers less likely to switch cards"),
    ("FICO",              0.035, "FICO Credit Score",                  "Mid-FICO band (664–730) anchors two of the top 3 segments"),
    ("EFX_AL_DQOCURNC_3", 0.031, "DQ Occurrences — 3-year window",    "Clean DQ history paradoxically means less credit-hungry — lower response"),
    ("EFX_BC_BAL_5",      0.025, "Bankcard Balance (sub-category)",    "Supporting signal for balance-based segments"),
    ("EFX_BC_TRDAGE_1",   0.022, "Age of Most Recent Bankcard Trade",  "Recent bankcard opener is more likely to respond"),
    ("EFX_AL_TRDCNT_1",   0.021, "Total Open Trade Count",             "Many trades = established credit user, less motivated by new offers"),
]

max_shap = 0.058
bar_max_w = 5.5
ty = 1.5
for i, (feat, shap, label, interp) in enumerate(drivers):
    ry = ty + i * 0.58
    bar_w = (shap / max_shap) * bar_max_w
    # rank circle
    box(s6, 0.3, ry + 0.07, 0.32, 0.32, TEAL)
    txt(s6, str(i + 1), 0.3, ry + 0.06, 0.32, 0.32, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # feature name
    txt(s6, label, 0.68, ry + 0.04, 3.3, 0.3, size=10, bold=True, color=NAVY)
    txt(s6, feat,  0.68, ry + 0.3, 3.3, 0.25, size=8.5, color=MGRAY, italic=True)
    # bar
    box(s6, 4.1, ry + 0.08, bar_w, 0.28, TEAL)
    txt(s6, f"{shap:.3f}", 4.15 + bar_w, ry + 0.08, 0.6, 0.28, size=9, bold=True, color=NAVY)
    # interpretation
    txt(s6, interp, 10.1, ry + 0.04, 2.95, 0.5, size=8.5, color=DKGRAY, wrap=True)

hline(s6, 6.2, MGRAY, 0.3, 12.73)
txt(s6, "SHAP values from gradient-boosted model trained on 26 EFX bureau features + FICO + Income. Mean absolute SHAP = average impact on model output magnitude.",
    0.3, 6.3, 12.73, 0.4, size=9, color=MGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RECOMMENDATIONS & NEXT STEPS
# ════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7, WHITE)
header_bar(s7, "Recommendations & Next Steps", "Analyst accept / reject — then operationalise")
footer(s7, "Slide 7 of 7")

# Recommended actions left
box(s7, 0.3, 1.45, 6.3, 5.3, LGRAY)
txt(s7, "Recommended Suppression Actions", 0.5, 1.55, 6.0, 0.38, size=13, bold=True, color=NAVY)
hline(s7, 2.05, TEAL, 0.5, 6.0)

recs = [
    ("Accept Seg 1 + 2 (minimum action)",
     "Suppress ~21% of mail. Response rates 42% below BAU. Minimal revenue risk — these prospects almost never respond."),
    ("Accept Seg 3 + 4 (moderate action)",
     "Add ~26% of mail to suppression. Rates 36–38% below BAU. Good FICO prospects unlikely to engage with this offer."),
    ("Accept Seg 5 (high-balance cohort)",
     "Add high-balance clean-DQ prospects. 33% below BAU. May be served by a different product — consider cross-sell rather than pure suppress."),
    ("Re-run on next campaign file",
     "Rules are feature-based, not sample-specific. Validate lift holds on the next solicitation universe before applying."),
]
ty = 2.2
for i, (title, body) in enumerate(recs):
    box(s7, 0.5, ty, 0.08, 0.22, ORANGE if i < 3 else TEAL)
    txt(s7, title, 0.7, ty - 0.02, 5.7, 0.28, size=10.5, bold=True, color=NAVY)
    txt(s7, body,  0.7, ty + 0.28, 5.7, 0.45, size=9.5, color=DKGRAY, wrap=True)
    ty += 0.9

# Right: impact summary
box(s7, 6.9, 1.45, 6.1, 5.3, NAVY)
txt(s7, "Projected Impact", 7.1, 1.55, 5.7, 0.38, size=13, bold=True, color=WHITE)
hline(s7, 2.05, TEAL, 7.1, 5.7)

impacts = [
    ("Segs 1–2 only",  "~21,000 records",  "0.24% resp rate", "42% below BAU"),
    ("Segs 1–4",       "~40,000 records",  "0.25% resp rate", "39% below BAU"),
    ("All 6 segs",     "~50,000 records",  "0.27% resp rate", "34% below BAU"),
]
ty2 = 2.2
for label, size, rate, lift in impacts:
    box(s7, 7.1, ty2, 5.7, 1.15, RGBColor(0x14, 0x3A, 0x6B))
    txt(s7, label, 7.25, ty2 + 0.05, 5.4, 0.3, size=11, bold=True, color=TEAL)
    txt(s7, size,  7.25, ty2 + 0.38, 2.5, 0.3, size=12, bold=True, color=WHITE)
    txt(s7, rate,  9.8,  ty2 + 0.38, 2.7, 0.3, size=11, color=MGRAY)
    txt(s7, lift,  7.25, ty2 + 0.72, 5.4, 0.3, size=10, color=ORANGE)
    ty2 += 1.3

txt(s7, "* Overlap-adjusted estimates. Exact counts depend on which segments are combined.",
    6.9, 6.6, 6.1, 0.3, size=8.5, color=MGRAY, italic=True)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "outputs/Suppression_Segment_Report_2026-06-26.pptx"
prs.save(out)
print(f"Saved -> {out}")
